"""
fill_competitor_summary.py — 競合比較サマリースライドをPPTXネイティブテーブルで生成

テンプレート構造:
  - Title 1            (PLACEHOLDER): top placeholder (brand 別)
  - Text Placeholder 2 (PLACEHOLDER): subtitle placeholder (brand 別)
  - Content Area       (AUTO_SHAPE):  削除してネイティブテーブルに置換
  - Source / Source 3  (TEXT_BOX):    出典 (stella=Source / roleup=Source 3)

方式: 横型比較テーブル（行=比較項目、列=企業）をネイティブテーブルで生成。
      target_company を指定した場合のみ、対象会社の列をハイライト。
      target_company が未指定なら competitors[] のみで全社フラット表示。
      3〜5競合の可変列数に対応。

Usage:
  python fill_competitor_summary.py \
    --data /home/claude/competitor_summary_data.json \
    --brand stellar_aiz | roleup \
    --output /mnt/user-data/outputs/CompetitorSummary_output.pptx
"""

import argparse
import copy
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "competitor-summary-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_CONTENT_AREA = "Content Area"
SHAPE_SOURCE_STELLA = "Source"      # stella template
SHAPE_SOURCE_ROLEUP = "Source 3"    # roleup (cp-derived) template

# ── レイアウト定数 (defaults; reassigned by _apply_theme) ──
CONTENT_LEFT = Inches(0.41)
CONTENT_TOP = Inches(1.50)
CONTENT_WIDTH = Inches(12.52)
CONTENT_BOTTOM = Inches(6.90)
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP

# 比較項目列（一番左）の幅比率
ITEM_COL_RATIO = 0.14

# ── 色定数 (defaults; reassigned by _apply_theme) ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_HEADER_BG = RGBColor(0xF0, 0xF0, 0xF0)
COLOR_TARGET_BG = RGBColor(0xFF, 0xF4, 0xC2)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_EVEN_ROW = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_BORDER = RGBColor(0xD0, 0xD0, 0xD0)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)

# ── フォント (defaults) ──
FONT_NAME_JP = "Meiryo UI"
FONT_NAME_LATIN = "Arial"

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_BOTTOM, CONTENT_HEIGHT
    global COLOR_TEXT, COLOR_HEADER_BG, COLOR_TARGET_BG, COLOR_EVEN_ROW
    global FONT_NAME_JP, FONT_NAME_LATIN

    _THEME = theme

    CONTENT_LEFT = theme.layout("content_left_in")
    CONTENT_TOP = theme.layout("content_top_in")
    CONTENT_WIDTH = theme.layout("content_width_in")
    CONTENT_BOTTOM = theme.layout("content_bottom_in")
    CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP

    COLOR_TEXT = theme.color("text")

    if theme.id == "roleup":
        COLOR_HEADER_BG = theme.color("header_bg")     # 薄褐色
        COLOR_TARGET_BG = theme.color("highlight_target")
        COLOR_EVEN_ROW = _hex_to_rgbcolor(theme.hex("label_bg"))
        FONT_NAME_JP = theme.font_ea
        FONT_NAME_LATIN = theme.font_latin
    # else stella: keep V1 hardcoded values for regression-zero.


def _hex_to_rgbcolor(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _silent_remove_shape(slide, shape_name: str) -> None:
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def remove_shape(slide, name):
    shape = find_shape(slide, name)
    if shape is not None:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)
        print(f"  ✓ Shape '{name}' removed")


def set_textbox_text(shape, text):
    """TextBoxのテキストを上書き（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def get_font_sizes(num_competitors):
    """競合社数に応じてフォントサイズを動的決定（対象+5社まで）。

    roleup: C4 [22,14,12,10,6] 制約のため固定 12/10/12pt。
    stella: 既存の動的調整 (14/13/12pt 系)。
    """
    if _THEME is not None and _THEME.id == "roleup":
        return {"header": 12, "body": 10, "item": 12}

    num_companies = num_competitors + 1
    if num_companies <= 4:
        return {"header": 14, "body": 13, "item": 14}
    elif num_companies == 5:
        return {"header": 13, "body": 12, "item": 13}
    else:
        return {"header": 12, "body": 11, "item": 12}


def apply_cell_style(cell, text, *,
                     font_size=10, bold=False,
                     bg_color=None, text_color=None,
                     align="left", v_align="top", is_multiline=False):
    """セルにテキストとスタイルを適用"""
    if bg_color is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

    if v_align == "middle":
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif v_align == "bottom":
        cell.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        cell.vertical_anchor = MSO_ANCHOR.TOP

    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)

    tf = cell.text_frame
    tf.word_wrap = True

    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    if isinstance(text, list):
        lines = [str(line) for line in text]
    else:
        lines = [str(text)] if text else [""]

    align_map = {"left": "l", "center": "ctr", "right": "r"}
    algn_code = align_map.get(align, "l")

    text_rgb = text_color if text_color is not None else COLOR_TEXT

    for line in lines:
        p_elem = etree.SubElement(txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        pPr.set("algn", algn_code)

        r_elem = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(font_size * 100),
            "b": "1" if bold else "0",
        })

        solidFill = etree.SubElement(rPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", f"{text_rgb[0]:02X}{text_rgb[1]:02X}{text_rgb[2]:02X}")

        latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", FONT_NAME_LATIN)
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", FONT_NAME_JP)
        cs = etree.SubElement(rPr, qn("a:cs"))
        cs.set("typeface", FONT_NAME_JP)

        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = line


def build_table(slide, data):
    """競合比較テーブルを動的に構築"""
    target = data.get("target_company") or None
    has_target = bool(target and target.get("name"))
    competitors = data["competitors"]
    items = data["comparison_items"]

    companies = ([target] + list(competitors)) if has_target else list(competitors)

    n_companies = len(companies)
    n_rows = len(items) + 1
    n_cols = 1 + n_companies

    # Content Areaを削除 (両 brand のテンプレに存在する想定)
    content_shape = find_shape(slide, SHAPE_CONTENT_AREA)
    if content_shape is not None:
        tbl_left = content_shape.left
        tbl_top = content_shape.top
        tbl_width = content_shape.width
        tbl_height = content_shape.height
        sp_tree = slide.shapes._spTree
        sp_tree.remove(content_shape._element)
    else:
        tbl_left = CONTENT_LEFT
        tbl_top = CONTENT_TOP
        tbl_width = CONTENT_WIDTH
        tbl_height = CONTENT_HEIGHT

    item_col_w = int(tbl_width * ITEM_COL_RATIO)
    remaining_w = tbl_width - item_col_w
    company_col_w = remaining_w // n_companies

    header_row_h = Inches(0.40)
    available_h = tbl_height - header_row_h
    data_row_h = int(available_h / len(items))

    fs = get_font_sizes(n_companies - 1)

    print(f"  ✓ テーブル: {n_rows}行 × {n_cols}列")
    if has_target:
        print(f"    対象会社+競合={n_companies}社, 列幅(企業): {company_col_w/914400:.2f}in")
    else:
        print(f"    全社フラット表示={n_companies}社, 列幅(企業): {company_col_w/914400:.2f}in")
    print(f"    データ行高: {data_row_h/914400:.2f}in, フォント: body={fs['body']}pt")

    table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table_shape.name = "CompetitorSummaryTable"
    table = table_shape.table

    table.columns[0].width = item_col_w
    for c in range(1, n_cols):
        table.columns[c].width = company_col_w

    tbl_xml = table_shape._element.find('.//' + qn('a:tbl'))
    for i, tr in enumerate(tbl_xml.findall(qn('a:tr'))):
        if i == 0:
            tr.set('h', str(header_row_h))
        else:
            tr.set('h', str(data_row_h))

    tbl_elem = table_shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    # ── ヘッダー行 ──
    apply_cell_style(
        table.cell(0, 0), "",
        font_size=fs["header"], bold=True,
        bg_color=COLOR_HEADER_BG,
        align="center", v_align="middle",
    )

    for c_idx, comp in enumerate(companies):
        is_target_col = has_target and c_idx == 0
        apply_cell_style(
            table.cell(0, 1 + c_idx),
            comp.get("name", f"競合{c_idx+1}"),
            font_size=fs["header"], bold=True,
            bg_color=COLOR_TARGET_BG if is_target_col else COLOR_HEADER_BG,
            align="center", v_align="middle",
        )

    # ── データ行 ──
    for r_idx, item in enumerate(items):
        item_key = item.get("key")
        item_label = item.get("label", item_key)
        row = r_idx + 1

        row_bg = COLOR_EVEN_ROW if (r_idx % 2 == 0) else COLOR_WHITE

        apply_cell_style(
            table.cell(row, 0), item_label,
            font_size=fs["item"], bold=True,
            bg_color=COLOR_HEADER_BG,
            align="left", v_align="middle",
        )

        for c_idx, comp in enumerate(companies):
            is_target_col = has_target and c_idx == 0
            val = comp.get(item_key, "")
            apply_cell_style(
                table.cell(row, 1 + c_idx), val,
                font_size=fs["body"],
                bold=is_target_col,
                bg_color=COLOR_TARGET_BG if is_target_col else row_bg,
                align="left", v_align="top",
            )

        print(f"    行{row}: {item_label}")

    print(f"  ✓ テーブル生成完了")


MAIN_MESSAGE_MAX = 65
CELL_VALUE_MAX = 30
COMPETITORS_MIN = 2
COMPETITORS_MAX = 5


def _validate_input(data):
    """main_message ≤65字、competitors=2〜5、各セル値 ≤30字。"""
    main_message = data.get("main_message", "")
    if len(main_message) > MAIN_MESSAGE_MAX:
        raise ValueError(
            f"main_message は {MAIN_MESSAGE_MAX} 字以内（受領: {len(main_message)}）: {main_message[:80]}..."
        )
    target = data.get("target_company") or {}
    has_target = bool(target.get("name"))
    competitors = data.get("competitors", [])
    if not (COMPETITORS_MIN <= len(competitors) <= COMPETITORS_MAX):
        raise ValueError(
            f"competitors の要素数は {COMPETITORS_MIN}〜{COMPETITORS_MAX} の範囲である必要があります"
            f"（受領: {len(competitors)}）"
        )
    comparison_items = data.get("comparison_items", [])
    keys = [item.get("key") for item in comparison_items if item.get("key")]
    for k in keys:
        if k in ("name",):
            continue
        if has_target:
            v = target.get(k, "")
            if isinstance(v, str) and len(v) > CELL_VALUE_MAX:
                raise ValueError(
                    f"target_company.{k} は {CELL_VALUE_MAX} 字以内（受領: {len(v)}）: {v}"
                )
        for i, c in enumerate(competitors):
            cv = c.get(k, "")
            if isinstance(cv, str) and len(cv) > CELL_VALUE_MAX:
                raise ValueError(
                    f"competitors[{i}].{k} は {CELL_VALUE_MAX} 字以内（受領: {len(cv)}）: {cv}"
                )


def main():
    parser = argparse.ArgumentParser(description="競合比較サマリー PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True, help="JSONデータファイルのパス")
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand.",
    )
    parser.add_argument("--output", required=True, help="出力PPTXファイルのパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "competitor-summary")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "target_company", "competitors", "comparison_items"],
        allowed_top=[
            "main_message", "chart_title", "source",
            "target_company", "competitors", "comparison_items",
            "title", "subtitle",
        ],
        per_item_required={
            "competitors": ["name"],
            "comparison_items": ["key", "label"],
        },
        skill_name=SKILL_ID,
    )

    require_source(data, theme, skill_id=SKILL_ID)
    _validate_input(data)
    n_comp = len(data.get("competitors", []))
    target_name = (data.get("target_company") or {}).get("name") or "（指定なし・強調なしモード）"
    print(f"  データ読み込み: 対象={target_name}, 競合={n_comp}社")

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top / subtitle placeholder (brand 別)
    top_text = resolve_top_text(data, theme)
    sub_text = resolve_subtitle_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE),
                     sub_text or data.get("chart_title", "競合比較"))

    # Roleup: silently remove brown guide rectangles carried by cp-derived template.
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # 3. テーブル構築
    build_table(slide, data)

    # 4. 出典 (stella="Source", roleup="Source 3")
    source_text = data.get("source", "出典：各社HP、IR資料、東京商工リサーチ等")
    source_shape_name = SHAPE_SOURCE_ROLEUP if theme.id == "roleup" else SHAPE_SOURCE_STELLA
    source_shape = find_shape(slide, source_shape_name)
    if source_shape is not None:
        set_textbox_text(source_shape, source_text)
        print(f"  ✓ Source ({source_shape_name}): {source_text[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
