"""
fill_financial_benchmark.py — 財務ベンチマーク比較スライドをPPTXネイティブオブジェクトで生成

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。
チャートは MSO_SHAPE.RECTANGLE で手動描画 (chart object 不在のため C10/C12 適用外)。

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 2×3グリッド: 6つの小型バーチャート（1指標/1チャート）
  - 下部: 出典

各小型チャート:
  - タイトル（指標名 + 単位）
  - 水平バーチャート（MSO_SHAPE.RECTANGLE で手動描画）
  - 対象会社は強調色でハイライト、その他は通常色
  - 企業名（左）+ 値ラベル（バーの右端）

Usage:
  python fill_financial_benchmark.py --brand stellar_aiz \
    --data /home/claude/financial_benchmark_data.json \
    --output /mnt/user-data/outputs/FinancialBenchmark_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "financial-benchmark-pptx"

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Layout / Style globals (defaults reassigned by _apply_theme) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE = "Source 3"

GRID_LEFT = Inches(0.41)
GRID_TOP = Inches(1.50)
GRID_WIDTH = Inches(12.51)
GRID_HEIGHT = Inches(5.40)
CELL_GAP_X = Inches(0.15)
CELL_GAP_Y = Inches(0.20)
LABEL_COL_W = Inches(1.15)
VALUE_COL_W = Inches(0.85)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.00)
SOURCE_W = Inches(12.50)
SOURCE_H = Inches(0.25)

N_COLS = 3
N_ROWS = 2

# 各セル内部のレイアウト
CHART_TITLE_H = Inches(0.35)
BAR_AREA_TOP_MARGIN = Inches(0.10)
BAR_AREA_BOTTOM_MARGIN = Inches(0.05)
BAR_AREA_LEFT_MARGIN = Inches(0.10)
BAR_AREA_RIGHT_MARGIN = Inches(0.10)

# Defaults (stella V1 colors; reassigned in main() via _apply_theme).
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_CELL_BG = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_CELL_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
COLOR_TITLE_UNDERLINE = RGBColor(0x33, 0x33, 0x33)
COLOR_BAR_DEFAULT = RGBColor(0x4E, 0x79, 0xA7)
COLOR_BAR_TARGET = RGBColor(0xE1, 0x57, 0x59)
COLOR_BAR_NEGATIVE = RGBColor(0xB8, 0x3A, 0x3A)
COLOR_BAR_TARGET_NEG = RGBColor(0x8B, 0x2C, 0x2E)
COLOR_NA = RGBColor(0x99, 0x99, 0x99)

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_CHART_TITLE = Pt(11)
FONT_SIZE_COMPANY = Pt(10)
FONT_SIZE_VALUE = Pt(10)
FONT_SIZE_SOURCE = Pt(10)

# Theme module-global; populated in main() via _apply_theme(theme).
_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global GRID_LEFT, GRID_TOP, GRID_WIDTH, GRID_HEIGHT, CELL_GAP_X, CELL_GAP_Y
    global LABEL_COL_W, VALUE_COL_W, SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE, COLOR_CELL_BG, COLOR_CELL_BORDER, COLOR_TITLE_UNDERLINE
    global COLOR_BAR_DEFAULT, COLOR_BAR_TARGET, COLOR_BAR_NEGATIVE, COLOR_BAR_TARGET_NEG
    global FONT_NAME_JP, FONT_SIZE_CHART_TITLE, FONT_SIZE_COMPANY, FONT_SIZE_VALUE, FONT_SIZE_SOURCE

    _THEME = theme

    GRID_LEFT = theme.layout("grid_left_in")
    GRID_TOP = theme.layout("grid_top_in")
    GRID_WIDTH = theme.layout("grid_width_in")
    GRID_HEIGHT = theme.layout("grid_height_in")
    CELL_GAP_X = theme.layout("cell_gap_x_in")
    CELL_GAP_Y = theme.layout("cell_gap_y_in")
    LABEL_COL_W = theme.layout("label_col_w_in")
    VALUE_COL_W = theme.layout("value_col_w_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")
    COLOR_TITLE_UNDERLINE = theme.color("text")
    COLOR_BAR_DEFAULT = theme.color("label_bar")
    COLOR_BAR_TARGET = theme.color("highlight_target")

    FONT_NAME_JP = theme.font_ea

    if theme.id == "stellar_aiz":
        # V1 hardcoded values for regression-zero.
        COLOR_CELL_BG = RGBColor(0xFA, 0xFA, 0xFA)
        COLOR_CELL_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
        COLOR_BAR_NEGATIVE = RGBColor(0xB8, 0x3A, 0x3A)
        COLOR_BAR_TARGET_NEG = RGBColor(0x8B, 0x2C, 0x2E)
        FONT_SIZE_CHART_TITLE = Pt(11)
        FONT_SIZE_COMPANY = Pt(10)
        FONT_SIZE_VALUE = Pt(10)
        FONT_SIZE_SOURCE = Pt(10)
    else:
        # Roleup: brand 一貫性で cell BG/border は brand 色、negative bar は
        # accent_op_margin_line で代用 (ユーザー判断 2026-05-06)。bold で対象
        # 識別を担保。
        COLOR_CELL_BG = theme.color("label_bg")
        COLOR_CELL_BORDER = theme.color("highlight_other")
        COLOR_BAR_NEGATIVE = theme.color("accent_op_margin_line")
        COLOR_BAR_TARGET_NEG = theme.color("accent_op_margin_line")
        # C4 allowed set: {22, 14, 12, 10, 6}
        body = theme.font_size_body_pt(skill_id=SKILL_ID)
        FONT_SIZE_CHART_TITLE = body
        FONT_SIZE_COMPANY = body
        FONT_SIZE_VALUE = body
        FONT_SIZE_SOURCE = theme.pt("font_size_source_pt")


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


# ──────────────────────────────────────────────
# Utility
# ──────────────────────────────────────────────
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                 font_name=None):
    # font_name / color default to current module globals (post-_apply_theme).
    # Late binding via None sentinel avoids capturing pre-theme defaults.
    if font_name is None:
        font_name = FONT_NAME_JP
    if color is None:
        color = COLOR_TEXT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return tb


def format_value(val, unit="", decimals=1, show_sign=False):
    """数値を整形して文字列に変換"""
    if val is None:
        return "—"

    if decimals == 0:
        num_str = f"{val:,.0f}"
    else:
        num_str = f"{val:,.{decimals}f}"

    if show_sign and val > 0 and not num_str.startswith("+"):
        num_str = "+" + num_str

    if unit:
        return f"{num_str}{unit}"
    return num_str


# ──────────────────────────────────────────────
# Single Bar Chart Cell
# ──────────────────────────────────────────────
def draw_bar_chart_cell(slide, metric, companies, target_company,
                         left, top, width, height):
    """1つの指標に対する水平バーチャートを描画する"""
    # セル背景
    cell_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    cell_bg.fill.solid()
    cell_bg.fill.fore_color.rgb = COLOR_CELL_BG
    cell_bg.line.color.rgb = COLOR_CELL_BORDER
    cell_bg.line.width = Pt(0.5)
    cell_bg.shadow.inherit = False
    cell_bg.text_frame.text = ""

    # 指標タイトル
    metric_name = metric.get("name", "指標")
    unit = metric.get("unit", "")
    title_text = f"{metric_name}({unit})" if unit else metric_name

    add_text_box(
        slide, title_text,
        left + Inches(0.10), top + Inches(0.05),
        width - Inches(0.20), CHART_TITLE_H,
        FONT_SIZE_CHART_TITLE, bold=True, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
    )

    # タイトル下の区切り線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.10), top + CHART_TITLE_H + Inches(0.03),
        width - Inches(0.20), Emu(int(Inches(0.01))),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TITLE_UNDERLINE
    line.line.fill.background()

    bar_area_top = top + CHART_TITLE_H + BAR_AREA_TOP_MARGIN + Inches(0.05)
    bar_area_h = height - CHART_TITLE_H - BAR_AREA_TOP_MARGIN - BAR_AREA_BOTTOM_MARGIN - Inches(0.05)
    bar_area_left = left + BAR_AREA_LEFT_MARGIN
    bar_area_right = left + width - BAR_AREA_RIGHT_MARGIN

    bar_start_x = bar_area_left + LABEL_COL_W
    bar_max_w = bar_area_right - bar_start_x - VALUE_COL_W

    values_dict = metric.get("values", {})
    decimals = metric.get("decimals", 1)
    show_sign = metric.get("show_sign", False)

    rows = []
    for comp in companies:
        name = comp["name"]
        val = values_dict.get(name)
        rows.append({"name": name, "value": val})

    sort_order = metric.get("sort", "keep")
    if sort_order == "desc":
        rows.sort(key=lambda r: (r["value"] if r["value"] is not None else float("-inf")), reverse=True)
    elif sort_order == "asc":
        rows.sort(key=lambda r: (r["value"] if r["value"] is not None else float("inf")))

    vals_numeric = [r["value"] for r in rows if r["value"] is not None]
    if not vals_numeric:
        return
    max_abs = max(abs(v) for v in vals_numeric) or 1
    has_negative = any(v < 0 for v in vals_numeric)

    n_rows = len(rows)
    row_h = bar_area_h // n_rows
    bar_h = Emu(int(row_h * 0.60))

    if has_negative:
        zero_x = bar_start_x + bar_max_w // 2
        half_bar_max_w = bar_max_w // 2
    else:
        zero_x = bar_start_x
        half_bar_max_w = bar_max_w

    for i, row in enumerate(rows):
        name = row["name"]
        val = row["value"]
        is_target = (name == target_company) if target_company else False

        row_y = bar_area_top + row_h * i
        cell_y_center = row_y + row_h // 2
        bar_y = cell_y_center - bar_h // 2

        # 企業名ラベル
        label_x = bar_area_left
        label_tb = slide.shapes.add_textbox(
            label_x, row_y, LABEL_COL_W, row_h,
        )
        ltf = label_tb.text_frame
        ltf.word_wrap = False
        ltf.margin_left = 0; ltf.margin_right = Inches(0.05)
        ltf.margin_top = 0; ltf.margin_bottom = 0
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.RIGHT
        lrun = lp.add_run()
        lrun.text = name
        lrun.font.size = FONT_SIZE_COMPANY
        lrun.font.bold = True if is_target else False
        lrun.font.name = FONT_NAME_JP
        lrun.font.color.rgb = COLOR_BAR_TARGET if is_target else COLOR_TEXT

        # バー描画
        if val is None:
            na_tb = slide.shapes.add_textbox(
                bar_start_x, row_y, bar_max_w, row_h,
            )
            ntf = na_tb.text_frame
            ntf.margin_left = Inches(0.05); ntf.margin_right = 0
            ntf.margin_top = 0; ntf.margin_bottom = 0
            ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
            np = ntf.paragraphs[0]
            np.alignment = PP_ALIGN.LEFT
            nrun = np.add_run()
            nrun.text = "N/A"
            nrun.font.size = FONT_SIZE_VALUE
            nrun.font.name = FONT_NAME_JP
            nrun.font.color.rgb = COLOR_NA
            nrun.font.italic = True
            continue

        bar_length = Emu(int(half_bar_max_w * abs(val) / max_abs))

        if val >= 0:
            bar_x = zero_x
            bar_color = COLOR_BAR_TARGET if is_target else COLOR_BAR_DEFAULT
        else:
            bar_x = zero_x - bar_length
            bar_color = COLOR_BAR_TARGET_NEG if is_target else COLOR_BAR_NEGATIVE

        if bar_length > Emu(100):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                bar_x, bar_y, bar_length, bar_h,
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = bar_color
            bar.line.fill.background()
            bar.shadow.inherit = False
            bar.text_frame.text = ""

        val_text = format_value(val, unit="", decimals=decimals, show_sign=show_sign)

        if val >= 0:
            vlabel_x = bar_x + bar_length + Inches(0.04)
            vlabel_w = VALUE_COL_W
            vlabel_align = PP_ALIGN.LEFT
        else:
            vlabel_x = bar_x - VALUE_COL_W - Inches(0.04)
            vlabel_w = VALUE_COL_W
            vlabel_align = PP_ALIGN.RIGHT

        if vlabel_x + vlabel_w > bar_area_right:
            vlabel_x = bar_area_right - vlabel_w
        if vlabel_x < bar_area_left + LABEL_COL_W:
            vlabel_x = bar_area_left + LABEL_COL_W

        vlabel_tb = slide.shapes.add_textbox(
            vlabel_x, row_y, vlabel_w, row_h,
        )
        vtf = vlabel_tb.text_frame
        vtf.word_wrap = False
        vtf.margin_left = 0; vtf.margin_right = 0
        vtf.margin_top = 0; vtf.margin_bottom = 0
        vtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        vp = vtf.paragraphs[0]
        vp.alignment = vlabel_align
        vrun = vp.add_run()
        vrun.text = val_text
        vrun.font.size = FONT_SIZE_VALUE
        vrun.font.bold = True if is_target else False
        vrun.font.name = FONT_NAME_JP
        vrun.font.color.rgb = COLOR_BAR_TARGET if is_target else COLOR_TEXT

    print(f"    ✓ {metric_name}: {len(rows)}社")


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "financial-benchmark")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "companies", "metrics"],
        allowed_top=[
            "main_message", "chart_title", "source",
            "target_company", "companies", "metrics",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    require_source(data, theme, skill_id=SKILL_ID)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    top_text = resolve_top_text(data, theme)
    sub_text = resolve_subtitle_text(data, theme) or "財務ベンチマーク"
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:40]}")
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:40]}")

    # Roleup: silently remove brown guide rectangles carried by template.
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    companies = data.get("companies", [])
    metrics = data.get("metrics", [])
    target_company = data.get("target_company")

    if not companies or not metrics:
        print("  ✗ ERROR: 'companies' and 'metrics' are required", file=sys.stderr)
        sys.exit(1)

    n_metrics = len(metrics)
    if n_metrics > N_ROWS * N_COLS:
        print(f"  ⚠ WARNING: {n_metrics} metrics > grid capacity ({N_ROWS * N_COLS}). Only first {N_ROWS * N_COLS} will be shown.", file=sys.stderr)
        metrics = metrics[: N_ROWS * N_COLS]
        n_metrics = len(metrics)

    cell_w = (GRID_WIDTH - CELL_GAP_X * (N_COLS - 1)) / N_COLS
    cell_h = (GRID_HEIGHT - CELL_GAP_Y * (N_ROWS - 1)) / N_ROWS

    print(f"\n  各指標のチャート生成:")
    for i, metric in enumerate(metrics):
        row = i // N_COLS
        col = i % N_COLS
        cell_x = GRID_LEFT + (cell_w + CELL_GAP_X) * col
        cell_y = GRID_TOP + (cell_h + CELL_GAP_Y) * row

        draw_bar_chart_cell(
            slide, metric, companies, target_company,
            cell_x, cell_y, cell_w, cell_h,
        )

    # 出典: roleup は Source 3 placeholder、stella は dynamic textbox。
    source = data.get("source", "")
    if source:
        if theme.id == "stellar_aiz":
            add_text_box(
                slide, source,
                SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            )
        else:
            source_shape = find_shape(slide, SHAPE_SOURCE)
            if source_shape is not None:
                set_textbox_text(source_shape, source)
            else:
                add_text_box(
                    slide, source,
                    SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                    FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
