"""
fill_business_overview.py — 事業セグメント概要スライドをPPTXネイティブオブジェクトで生成

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。

レイアウト (customer-profile-pptx と同一構造):
  - 左側: 事業の概要テーブル (key-value 形式, ブレットポイント)
  - 右側 (mode=revenue_chart): 棒+折れ線複合チャート + CAGR 注釈 + カスタム凡例
  - 右側 (mode=kpi_cards): 1〜6 個の KPI カードを 2 列グリッドに配置
  - 下部: 出典 (stella=動的 textbox / roleup=Source 3 placeholder)

Usage:
  python fill_business_overview.py --brand stellar_aiz \\
    --data {{WORK_DIR}}/business_overview_data.json \\
    --output {{OUTPUT_DIR}}/BusinessOverview_output.pptx
"""

import argparse
import copy
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import (  # noqa: E402
    resolve_top_text,
    resolve_subtitle_text,
    require_source,
    apply_line_spacing,
    format_fiscal_period,
)

SKILL_ID = "business-overview-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names (template-structure invariants, brand-agnostic) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

# Defaults (stella). Reassigned in _apply_theme(theme).
SHAPE_SOURCE = "Source"
PANEL_Y = Inches(1.50)
LEFT_X = Inches(0.41)
LEFT_W = Inches(5.80)
RIGHT_X = Inches(6.50)
RIGHT_W = Inches(6.40)
CHART_H = Inches(4.80)
LABEL_COL_W = Inches(1.60)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(8.00)
SOURCE_H = Inches(0.30)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_BAR = RGBColor(0x4E, 0x79, 0xA7)
COLOR_LINE = RGBColor(0x00, 0x33, 0x66)
COLOR_CAGR_ARROW = RGBColor(0x33, 0x33, 0x33)
COLOR_KPI_CARD_BG = RGBColor(0xF7, 0xF7, 0xF7)
COLOR_KPI_CARD_BORDER = RGBColor(0xD0, 0xD0, 0xD0)
COLOR_KPI_VALUE = RGBColor(0x4E, 0x79, 0xA7)
COLOR_KPI_SUB = RGBColor(0x66, 0x66, 0x66)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_SUBTITLE = RGBColor(0x33, 0x33, 0x33)

TEXT_HEX = "333333"
ACCENT_REVENUE_BAR_HEX = "4E79A7"
ACCENT_OP_MARGIN_LINE_HEX = "003366"

FONT_NAME_JP = "Meiryo UI"

# Stella defaults; reassigned in _apply_theme(theme).
SECTION_TITLE_PT = 14
SECTION_TITLE_BOLD = True
SECTION_TITLE_ALIGN = PP_ALIGN.CENTER
SECTION_TITLE_UNDERLINE = True
LABEL_PT = 14
VALUE_PT = 14
KPI_NAME_PT = 12
KPI_VALUE_PT = 28
KPI_SUB_PT = 10
DATA_LABEL_SZ = "1200"  # 12pt (stella)
AXIS_FONT_SZ = "1100"   # 11pt (stella)
UNIT_FONT_PT = 12
LEGEND_FONT_PT = 12
CAGR_FONT_PT = 16
SOURCE_FONT_PT = 10

MAIN_MESSAGE_LIMIT = 65

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global SHAPE_SOURCE
    global PANEL_Y, LEFT_X, LEFT_W, RIGHT_X, RIGHT_W, CHART_H, LABEL_COL_W
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_BAR, COLOR_LINE, COLOR_CAGR_ARROW
    global COLOR_KPI_CARD_BG, COLOR_KPI_CARD_BORDER, COLOR_KPI_VALUE, COLOR_KPI_SUB
    global COLOR_SOURCE, COLOR_SUBTITLE
    global TEXT_HEX, ACCENT_REVENUE_BAR_HEX, ACCENT_OP_MARGIN_LINE_HEX
    global FONT_NAME_JP
    global SECTION_TITLE_PT, SECTION_TITLE_BOLD, SECTION_TITLE_ALIGN, SECTION_TITLE_UNDERLINE
    global LABEL_PT, VALUE_PT, KPI_NAME_PT, KPI_VALUE_PT, KPI_SUB_PT
    global DATA_LABEL_SZ, AXIS_FONT_SZ
    global UNIT_FONT_PT, LEGEND_FONT_PT, CAGR_FONT_PT, SOURCE_FONT_PT

    _THEME = theme
    FONT_NAME_JP = theme.font_ea
    TEXT_HEX = theme.hex_no_hash("text")
    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    PANEL_Y = theme.layout("panel_y_in")
    LEFT_X = theme.layout("left_x_in")
    LEFT_W = theme.layout("left_w_in")
    RIGHT_X = theme.layout("right_x_in")
    RIGHT_W = theme.layout("right_w_in")
    CHART_H = theme.layout("chart_h_in")
    LABEL_COL_W = theme.layout("label_col_w_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        # stella V1 互換: 旧 hardcode 維持
        COLOR_BAR = RGBColor(0x4E, 0x79, 0xA7)
        COLOR_LINE = RGBColor(0x00, 0x33, 0x66)
        COLOR_CAGR_ARROW = RGBColor(0x33, 0x33, 0x33)
        COLOR_KPI_CARD_BG = RGBColor(0xF7, 0xF7, 0xF7)
        COLOR_KPI_CARD_BORDER = RGBColor(0xD0, 0xD0, 0xD0)
        COLOR_KPI_VALUE = RGBColor(0x4E, 0x79, 0xA7)
        COLOR_KPI_SUB = RGBColor(0x66, 0x66, 0x66)
        COLOR_SUBTITLE = COLOR_TEXT
        ACCENT_REVENUE_BAR_HEX = "4E79A7"
        ACCENT_OP_MARGIN_LINE_HEX = "003366"
        SECTION_TITLE_PT = 14
        SECTION_TITLE_BOLD = True
        SECTION_TITLE_ALIGN = PP_ALIGN.CENTER
        SECTION_TITLE_UNDERLINE = True
        LABEL_PT = 14
        VALUE_PT = 14
        KPI_NAME_PT = 12
        KPI_VALUE_PT = 28
        KPI_SUB_PT = 10
        DATA_LABEL_SZ = "1200"
        AXIS_FONT_SZ = "1100"
        UNIT_FONT_PT = 12
        LEGEND_FONT_PT = 12
        CAGR_FONT_PT = 16
        SOURCE_FONT_PT = 10
    else:
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        SHAPE_SOURCE = "Source 3"
        COLOR_BAR = theme.color("accent_revenue_bar")
        COLOR_LINE = theme.color("accent_op_margin_line")
        COLOR_CAGR_ARROW = theme.color("cagr_arrow")
        # KPI card: 茶系トーン
        COLOR_KPI_CARD_BG = theme.color("label_bg")           # #F2E8DD
        COLOR_KPI_CARD_BORDER = theme.color("highlight_other")  # #CDCECE
        COLOR_KPI_VALUE = theme.color("accent_revenue_bar")    # #7C4C2C
        COLOR_KPI_SUB = theme.color("source")                  # #3E3A39
        COLOR_SUBTITLE = theme.color("subtitle")
        ACCENT_REVENUE_BAR_HEX = theme.hex_no_hash("accent_revenue_bar")
        ACCENT_OP_MARGIN_LINE_HEX = theme.hex_no_hash("accent_op_margin_line")
        SECTION_TITLE_PT = theme.pt_value("font_size_subtitle_pt")  # 12
        SECTION_TITLE_BOLD = False
        SECTION_TITLE_ALIGN = PP_ALIGN.LEFT
        SECTION_TITLE_UNDERLINE = False  # roleup template の object 8 が下線役
        LABEL_PT = theme.pt_value("font_size_body_pt")  # 10
        VALUE_PT = theme.pt_value("font_size_body_pt")  # 10
        KPI_NAME_PT = theme.pt_value("font_size_subtitle_pt")  # 12 (KPI 名はサブ見出し扱い)
        # KPI 値は featured display number。C4 許容セットの最大 22pt (タイトル相当) を採用。
        KPI_VALUE_PT = theme.pt_value("font_size_title_pt")  # 22
        KPI_SUB_PT = theme.pt_value("font_size_body_pt")  # 10
        DATA_LABEL_SZ = "1000"  # 10pt
        AXIS_FONT_SZ = "1000"   # 10pt
        UNIT_FONT_PT = theme.pt_value("font_size_body_pt")  # 10
        LEGEND_FONT_PT = theme.pt_value("font_size_body_pt")  # 10
        CAGR_FONT_PT = theme.pt_value("font_size_key_message_pt")  # 14 (C4: 16 不可)
        SOURCE_FONT_PT = theme.pt_value("font_size_source_pt")  # 6


def _silent_remove_shape(slide, name):
    """find_shape の warning を出さずに削除を試みる(brand 別 shape 名フォールバック用)"""
    for shape in list(slide.shapes):
        if shape.name == name:
            slide.shapes._spTree.remove(shape._element)
            return True
    return False


def _validate(data):
    main_message = data.get("main_message", "")
    if not main_message:
        raise SystemExit("ERROR: main_message is required")
    if len(main_message) > MAIN_MESSAGE_LIMIT:
        raise SystemExit(
            f"ERROR: main_message exceeds {MAIN_MESSAGE_LIMIT} chars "
            f"(actual: {len(main_message)}). Shorten it.\n"
            f"  text: {main_message!r}"
        )
    if not data.get("parent_company"):
        raise SystemExit("ERROR: parent_company is required")
    if not data.get("segment_name"):
        raise SystemExit("ERROR: segment_name is required")
    overview = data.get("overview") or {}
    if not overview.get("items"):
        raise SystemExit("ERROR: overview.items is required")
    perf = data.get("performance") or {}
    mode = perf.get("mode")
    if mode not in ("revenue_chart", "kpi_cards"):
        raise SystemExit(
            f"ERROR: performance.mode must be 'revenue_chart' or 'kpi_cards' (got {mode!r})"
        )
    if mode == "revenue_chart" and not perf.get("data"):
        raise SystemExit("ERROR: performance.data is required for revenue_chart mode")
    if mode == "kpi_cards" and not perf.get("cards"):
        raise SystemExit("ERROR: performance.cards is required for kpi_cards mode")


def find_shape(slide, name, warn=True):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    if warn:
        print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def remove_shape(slide, name):
    shape = find_shape(slide, name, warn=False)
    if shape is not None:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)
        print(f"  ✓ Shape '{name}' removed")


def add_section_title(slide, text, left, top, width):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = SECTION_TITLE_ALIGN
    run = p.add_run()
    run.text = text
    run.font.size = Pt(SECTION_TITLE_PT)
    run.font.bold = SECTION_TITLE_BOLD
    run.font.color.rgb = COLOR_SUBTITLE
    run.font.name = FONT_NAME_JP

    if SECTION_TITLE_UNDERLINE:
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = COLOR_TEXT
        underline.line.fill.background()
    return txBox


def build_overview_table(slide, items, left, top, width):
    """事業の概要テーブルを構築（ブレットポイント形式、枠線なし）"""
    n_rows = len(items)
    n_cols = 2
    col0_w = LABEL_COL_W
    col1_w = width - col0_w

    # 行高: roleup は line_height_pt+4pt を最低として利用可能高さに均等分配。
    # stella は 0.35in 固定 (regression-zero)。
    if _THEME is not None and _THEME.line_height_pt() is not None:
        min_floor = Inches((_THEME.line_height_pt() + 4) / 72.0)
        section_title_h = Inches(0.40)
        bottom_margin = Inches(0.02)
        available_h = SOURCE_Y - PANEL_Y - section_title_h - bottom_margin
        total_min_h = min_floor * n_rows
        if total_min_h < available_h:
            extra_per_row = (available_h - total_min_h) / n_rows
            min_row_h = min_floor + extra_per_row
        else:
            min_row_h = min_floor
    else:
        min_row_h = Inches(0.35)
    table_h = min_row_h * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h)
    table = shape.table

    table.columns[0].width = col0_w
    table.columns[1].width = col1_w

    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '0', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    for tr in tbl_elem.findall(qn('a:tr')):
        tr.set('h', str(min_row_h))

    for r_idx, item in enumerate(items):
        label = item.get("label", "")
        value = item.get("value", "")
        bullet_label = f"•  {label}"
        _style_cell(table.cell(r_idx, 0), bullet_label, True, Pt(LABEL_PT))
        _style_cell(table.cell(r_idx, 1), value, False, Pt(VALUE_PT))

    print(f"  ✓ 事業概要テーブル: {n_rows}行")
    return shape


def _style_cell(cell, text, bold, font_size):
    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        txBody = etree.SubElement(tc, qn("a:txBody"))

    old_bodyPr = txBody.find(qn("a:bodyPr"))
    if old_bodyPr is not None:
        txBody.remove(old_bodyPr)
    bodyPr = etree.SubElement(txBody, qn("a:bodyPr"), attrib={
        "wrap": "square",
        "lIns": "0", "rIns": "0",
        "tIns": "27432", "bIns": "27432",
        "anchor": "t",
    })
    txBody.insert(0, bodyPr)

    if txBody.find(qn("a:lstStyle")) is None:
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        txBody.insert(1, lstStyle)

    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    p_elem = etree.SubElement(txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "l")

    if _THEME is not None and _THEME.line_height_pt() is not None:
        apply_line_spacing(pPr, _THEME)
    else:
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        etree.SubElement(lnSpc, qn("a:spcPct"), attrib={"val": "100000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "0"})
    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    etree.SubElement(spcAft, qn("a:spcPts"), attrib={"val": "0"})

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(font_size.pt * 100)),
        "b": "1" if bold else "0",
    })
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr"), attrib={"val": TEXT_HEX})
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = str(text)

    old_tcPr = tc.find(qn("a:tcPr"))
    if old_tcPr is not None:
        tc.remove(old_tcPr)
    tcPr = etree.SubElement(tc, qn("a:tcPr"), attrib={
        "marL": "45720", "marR": "18288",
        "marT": "27432", "marB": "27432",
        "anchor": "t",
    })
    for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = etree.SubElement(tcPr, qn(border_name), attrib={"w": "0", "cmpd": "sng"})
        etree.SubElement(ln, qn("a:noFill"))


def build_combo_chart(slide, perf_data, left, top, width, height):
    """PowerPointネイティブ複合チャート（棒＋折れ線）"""
    from pptx.chart.data import CategoryChartData

    data = perf_data["data"]
    bar_label = perf_data.get("bar_label", "セグメント売上高")
    line_label = perf_data.get("line_label", "セグメント営業利益率")

    chart_data = CategoryChartData()
    if _THEME is not None and _THEME.fiscal_period_format():
        fy_month = perf_data.get("fiscal_year_end_month", 12)
        chart_data.categories = [
            format_fiscal_period(int(d["year"]), fy_month, _THEME) for d in data
        ]
    else:
        chart_data.categories = [d["year"] for d in data]
    chart_data.add_series(bar_label, [d["revenue"] for d in data])
    chart_data.add_series(line_label, [d["op_margin"] for d in data])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data
    )
    chart = chart_frame.chart

    plotArea = chart._chartSpace.chart.plotArea
    barChart = plotArea.findall(qn('c:barChart'))[0]

    sers = barChart.findall(qn('c:ser'))
    line_ser_xml = copy.deepcopy(sers[1])
    barChart.remove(sers[1])

    lineChart = etree.SubElement(plotArea, qn('c:lineChart'))
    etree.SubElement(lineChart, qn('c:grouping'), attrib={'val': 'standard'})
    etree.SubElement(lineChart, qn('c:varyColors'), attrib={'val': '0'})
    lineChart.append(line_ser_xml)

    marker_xml = line_ser_xml.find(qn('c:marker'))
    if marker_xml is None:
        marker_xml = etree.SubElement(line_ser_xml, qn('c:marker'))
    symbol = marker_xml.find(qn('c:symbol'))
    if symbol is None:
        symbol = etree.SubElement(marker_xml, qn('c:symbol'))
    symbol.set('val', 'circle')
    sz = marker_xml.find(qn('c:size'))
    if sz is None:
        sz = etree.SubElement(marker_xml, qn('c:size'))
    sz.set('val', '9')

    valAx = plotArea.findall(qn('c:valAx'))

    sec_valAx_id = "2094734553"
    sec_catAx_id = "2094734554"

    etree.SubElement(lineChart, qn('c:axId'), attrib={'val': sec_catAx_id})
    etree.SubElement(lineChart, qn('c:axId'), attrib={'val': sec_valAx_id})

    sec_catAx = etree.SubElement(plotArea, qn('c:catAx'))
    etree.SubElement(sec_catAx, qn('c:axId'), attrib={'val': sec_catAx_id})
    scaling = etree.SubElement(sec_catAx, qn('c:scaling'))
    etree.SubElement(scaling, qn('c:orientation'), attrib={'val': 'minMax'})
    etree.SubElement(sec_catAx, qn('c:delete'), attrib={'val': '1'})
    etree.SubElement(sec_catAx, qn('c:axPos'), attrib={'val': 'b'})
    etree.SubElement(sec_catAx, qn('c:crossAx'), attrib={'val': sec_valAx_id})

    sec_valAx_elem = etree.SubElement(plotArea, qn('c:valAx'))
    etree.SubElement(sec_valAx_elem, qn('c:axId'), attrib={'val': sec_valAx_id})
    scaling2 = etree.SubElement(sec_valAx_elem, qn('c:scaling'))
    etree.SubElement(scaling2, qn('c:orientation'), attrib={'val': 'minMax'})
    etree.SubElement(sec_valAx_elem, qn('c:delete'), attrib={'val': '1'})
    etree.SubElement(sec_valAx_elem, qn('c:axPos'), attrib={'val': 'r'})
    etree.SubElement(sec_valAx_elem, qn('c:numFmt'), attrib={
        'formatCode': '0.0"%"', 'sourceLinked': '0'
    })
    etree.SubElement(sec_valAx_elem, qn('c:crossAx'), attrib={'val': sec_catAx_id})
    etree.SubElement(sec_valAx_elem, qn('c:crosses'), attrib={'val': 'max'})

    bar_ser = barChart.findall(qn('c:ser'))[0]
    spPr = bar_ser.find(qn('c:spPr'))
    if spPr is None:
        spPr = etree.SubElement(bar_ser, qn('c:spPr'))
    sf = spPr.find(qn('a:solidFill'))
    if sf is None:
        sf = etree.SubElement(spPr, qn('a:solidFill'))
    for child in list(sf):
        sf.remove(child)
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': ACCENT_REVENUE_BAR_HEX})

    line_spPr = line_ser_xml.find(qn('c:spPr'))
    if line_spPr is None:
        line_spPr = etree.SubElement(line_ser_xml, qn('c:spPr'))
    ln = line_spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(line_spPr, qn('a:ln'))
    ln.set('w', '19050')
    line_sf = ln.find(qn('a:solidFill'))
    if line_sf is None:
        line_sf = etree.SubElement(ln, qn('a:solidFill'))
    for child in list(line_sf):
        line_sf.remove(child)
    etree.SubElement(line_sf, qn('a:srgbClr'), attrib={'val': ACCENT_OP_MARGIN_LINE_HEX})

    marker_spPr = marker_xml.find(qn('c:spPr'))
    if marker_spPr is None:
        marker_spPr = etree.SubElement(marker_xml, qn('c:spPr'))
    m_sf = etree.SubElement(marker_spPr, qn('a:solidFill'))
    etree.SubElement(m_sf, qn('a:srgbClr'), attrib={'val': ACCENT_OP_MARGIN_LINE_HEX})

    add_data_labels_to_ser(bar_ser, position='outEnd', num_format='0.0', font_color=TEXT_HEX)
    add_data_labels_to_ser(line_ser_xml, position='t', num_format='0.0', font_color='FFFFFF')

    chart.has_legend = False

    for vax in plotArea.findall(qn('c:valAx')):
        for mg in vax.findall(qn('c:majorGridlines')):
            vax.remove(mg)
        for mg in vax.findall(qn('c:minorGridlines')):
            vax.remove(mg)
    for cax in plotArea.findall(qn('c:catAx')):
        for mg in cax.findall(qn('c:majorGridlines')):
            cax.remove(mg)

    gapWidth = barChart.find(qn('c:gapWidth'))
    if gapWidth is None:
        gapWidth = etree.SubElement(barChart, qn('c:gapWidth'))
    gapWidth.set('val', '108')

    for ax in plotArea.findall(qn('c:catAx')) + plotArea.findall(qn('c:valAx')):
        delete_elem = ax.find(qn('c:delete'))
        if delete_elem is not None and delete_elem.get('val') == '1':
            continue
        txPr = ax.find(qn('c:txPr'))
        if txPr is None:
            txPr = etree.SubElement(ax, qn('c:txPr'))
        bodyPr = txPr.find(qn('a:bodyPr'))
        if bodyPr is None:
            bodyPr = etree.SubElement(txPr, qn('a:bodyPr'))
            txPr.insert(0, bodyPr)
        bodyPr.set('rot', '-5400000')
        bodyPr.set('vert', 'horz')

        if txPr.find(qn('a:lstStyle')) is None:
            etree.SubElement(txPr, qn('a:lstStyle'))
        p = txPr.find(qn('a:p'))
        if p is None:
            p = etree.SubElement(txPr, qn('a:p'))
        pPr = p.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.SubElement(p, qn('a:pPr'))
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn('a:defRPr'), attrib={'sz': AXIS_FONT_SZ})
        else:
            defRPr.set('sz', AXIS_FONT_SZ)
        if defRPr.find(qn('a:latin')) is None:
            etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
        if defRPr.find(qn('a:ea')) is None:
            etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})

    for vax in valAx:
        del_elem = vax.find(qn('c:delete'))
        if del_elem is None:
            del_elem = etree.SubElement(vax, qn('c:delete'))
        del_elem.set('val', '1')

    chart.has_title = False

    plotArea_spPr = plotArea.find(qn('c:spPr'))
    if plotArea_spPr is None:
        plotArea_spPr = etree.SubElement(plotArea, qn('c:spPr'))
    etree.SubElement(plotArea_spPr, qn('a:noFill'))

    print(f"  ✓ 複合チャート: {len(data)}年分")
    return chart_frame


def add_data_labels_to_ser(ser_xml, position='outEnd', num_format='0.0', font_color=None):
    if font_color is None:
        font_color = TEXT_HEX
    dLbls = ser_xml.find(qn('c:dLbls'))
    if dLbls is None:
        dLbls = etree.SubElement(ser_xml, qn('c:dLbls'))
    for child in list(dLbls):
        dLbls.remove(child)

    etree.SubElement(dLbls, qn('c:numFmt'), attrib={
        'formatCode': num_format, 'sourceLinked': '0'
    })
    etree.SubElement(dLbls, qn('c:showLegendKey'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showVal'), attrib={'val': '1'})
    etree.SubElement(dLbls, qn('c:showCatName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showSerName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showPercent'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showBubbleSize'), attrib={'val': '0'})

    pos_map = {'outEnd': 'outEnd', 't': 't', 'ctr': 'ctr'}
    etree.SubElement(dLbls, qn('c:dLblPos'), attrib={
        'val': pos_map.get(position, 'outEnd')
    })

    txPr = etree.SubElement(dLbls, qn('c:txPr'))
    etree.SubElement(txPr, qn('a:bodyPr'))
    etree.SubElement(txPr, qn('a:lstStyle'))
    p = etree.SubElement(txPr, qn('a:p'))
    pPr = etree.SubElement(p, qn('a:pPr'))
    defRPr = etree.SubElement(pPr, qn('a:defRPr'), attrib={'sz': DATA_LABEL_SZ})
    etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
    etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})
    sf = etree.SubElement(defRPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': font_color})


def add_cagr_annotation(slide, perf_data, chart_left, chart_top, chart_width, chart_height):
    data = perf_data["data"]
    if len(data) < 2:
        return

    first_rev = data[0]["revenue"]
    last_rev = data[-1]["revenue"]
    max_rev = max(d["revenue"] for d in data)
    n_years = len(data) - 1
    n_cats = len(data)

    if first_rev > 0 and last_rev > 0 and n_years > 0:
        cagr = (last_rev / first_rev) ** (1.0 / n_years) - 1
        cagr_text = f"+{cagr*100:.1f}%" if cagr >= 0 else f"{cagr*100:.1f}%"
    else:
        cagr_text = "N/A"

    plot_left_margin = chart_width * 0.06
    plot_right_margin = chart_width * 0.04
    plot_top_margin = chart_height * 0.18
    plot_bottom_margin = chart_height * 0.14

    plot_left = chart_left + plot_left_margin
    plot_right = chart_left + chart_width - plot_right_margin
    plot_top = chart_top + plot_top_margin
    plot_bottom = chart_top + chart_height - plot_bottom_margin
    plot_w = plot_right - plot_left
    plot_h = plot_bottom - plot_top

    cat_width = plot_w / n_cats
    first_bar_cx = plot_left + 0.5 * cat_width
    last_bar_cx = plot_left + (n_cats - 0.5) * cat_width

    axis_max = max_rev * 1.20
    first_bar_top_y = plot_bottom - (first_rev / axis_max) * plot_h
    last_bar_top_y = plot_bottom - (last_rev / axis_max) * plot_h

    gap_above = Inches(1.20)
    arrow_start_x = int(first_bar_cx)
    arrow_start_y = int(first_bar_top_y - gap_above)
    arrow_end_x = int(last_bar_cx)
    arrow_end_y = int(last_bar_top_y - gap_above)

    connector = slide.shapes.add_connector(
        1,
        arrow_start_x, arrow_start_y,
        arrow_end_x, arrow_end_y
    )
    connector.line.color.rgb = COLOR_CAGR_ARROW
    connector.line.width = Pt(1.5)

    cxnSp = connector._element
    spPr = cxnSp.find(qn('p:spPr'))
    if spPr is None:
        spPr = cxnSp.find(qn('a:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    etree.SubElement(ln, qn('a:tailEnd'), attrib={
        'type': 'triangle', 'w': 'med', 'len': 'med'
    })

    text_w = Inches(1.30)
    text_h = Inches(0.40)
    mid_x = (arrow_start_x + arrow_end_x) / 2
    mid_y = (arrow_start_y + arrow_end_y) / 2
    text_x = int(mid_x - text_w / 2)
    text_y = int(mid_y - text_h / 2)

    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, text_x, text_y, text_w, text_h
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    oval.line.color.rgb = COLOR_CAGR_ARROW
    oval.line.width = Pt(1.0)

    tf = oval.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = cagr_text
    run.font.size = Pt(CAGR_FONT_PT)
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP

    print(f"  ✓ CAGR注釈: {cagr_text} ({data[0]['year']}→{data[-1]['year']})")


def add_unit_label(slide, text, left, top, width, height=None):
    if height is None:
        height = Inches(0.22)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(UNIT_FONT_PT)
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP


def add_custom_legend(slide, perf_data, left, top, width, max_right_emu=None):
    bar_label = perf_data.get("bar_label", "セグメント売上高")
    line_label = perf_data.get("line_label", "セグメント営業利益率")

    sq_size = Inches(0.14)
    bar_text_w = Inches(1.50)
    line_w = Inches(0.30)
    line_text_w = Inches(1.50)
    # Internal legend layout: marker + 0.06 gap + bar text + 0.10 gap + line + 0.06 gap + line text
    legend_w = sq_size + Inches(0.06) + bar_text_w + Inches(0.10) + line_w + Inches(0.06) + line_text_w
    legend_h = Inches(0.22)
    preferred_right = left + width
    if max_right_emu is not None and max_right_emu < preferred_right:
        legend_right = max_right_emu
    else:
        legend_right = preferred_right
    legend_x = legend_right - legend_w

    sq_y = top + (legend_h - sq_size) / 2

    bar_marker = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        legend_x, int(sq_y), sq_size, sq_size
    )
    bar_marker.fill.solid()
    bar_marker.fill.fore_color.rgb = COLOR_BAR
    bar_marker.line.fill.background()

    bar_text_x = legend_x + sq_size + Inches(0.06)
    txBox1 = slide.shapes.add_textbox(int(bar_text_x), top, bar_text_w, legend_h)
    tf1 = txBox1.text_frame
    tf1.word_wrap = False
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    run1 = p1.add_run()
    run1.text = bar_label
    run1.font.size = Pt(LEGEND_FONT_PT)
    run1.font.color.rgb = COLOR_TEXT
    run1.font.name = FONT_NAME_JP

    line_section_x = bar_text_x + bar_text_w + Inches(0.10)
    line_y = top + legend_h / 2
    connector = slide.shapes.add_connector(
        1, int(line_section_x), int(line_y),
        int(line_section_x + line_w), int(line_y)
    )
    connector.line.color.rgb = COLOR_LINE
    connector.line.width = Pt(1.5)

    circle_size = Inches(0.12)
    circle_x = line_section_x + (line_w - circle_size) / 2
    circle_y = top + (legend_h - circle_size) / 2
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(circle_x), int(circle_y), circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_LINE
    circle.line.fill.background()

    line_text_x = line_section_x + line_w + Inches(0.06)
    txBox2 = slide.shapes.add_textbox(int(line_text_x), top, line_text_w, legend_h)
    tf2 = txBox2.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = line_label
    run2.font.size = Pt(LEGEND_FONT_PT)
    run2.font.color.rgb = COLOR_TEXT
    run2.font.name = FONT_NAME_JP

    print(f"  ✓ カスタム凡例: {bar_label} / {line_label}")


def build_kpi_cards(slide, perf_data, left, top, width, height):
    """KPI カードグリッド（2 列 × 行可変）。各カード: KPI 名 / 値 / 補足。"""
    cards = perf_data["cards"]
    n = len(cards)
    if n < 1 or n > 6:
        raise SystemExit(f"ERROR: kpi_cards.cards must have 1-6 entries (got {n})")

    n_cols = 2
    n_rows = (n + 1) // 2
    gap = Inches(0.20)
    card_w = (width - gap) / n_cols
    card_h = (height - gap * (n_rows - 1)) / n_rows if n_rows > 1 else height

    for idx, card in enumerate(cards):
        row = idx // n_cols
        col = idx % n_cols
        cx = int(left + col * (card_w + gap))
        cy = int(top + row * (card_h + gap))

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, int(card_w), int(card_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_KPI_CARD_BG
        bg.line.color.rgb = COLOR_KPI_CARD_BORDER
        bg.line.width = Pt(0.75)
        bg.text_frame.text = ""

        name = card.get("name", "")
        value = card.get("value", "")
        unit = card.get("unit", "")
        sub = card.get("sub", "")

        # KPI 名 (上部)
        name_box = slide.shapes.add_textbox(
            cx + Inches(0.15), cy + Inches(0.10),
            int(card_w - Inches(0.30)), Inches(0.30)
        )
        tf_n = name_box.text_frame
        tf_n.word_wrap = True
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.LEFT
        run_n = p_n.add_run()
        run_n.text = name
        run_n.font.size = Pt(KPI_NAME_PT)
        run_n.font.bold = True
        run_n.font.color.rgb = COLOR_TEXT
        run_n.font.name = FONT_NAME_JP

        # KPI 値 (中央)
        value_h = Inches(0.70)
        value_box = slide.shapes.add_textbox(
            cx + Inches(0.15), cy + (card_h - value_h) / 2,
            int(card_w - Inches(0.30)), int(value_h)
        )
        tf_v = value_box.text_frame
        tf_v.word_wrap = False
        p_v = tf_v.paragraphs[0]
        p_v.alignment = PP_ALIGN.CENTER
        run_v = p_v.add_run()
        run_v.text = f"{value}{unit}" if unit else str(value)
        run_v.font.size = Pt(KPI_VALUE_PT)
        run_v.font.bold = True
        run_v.font.color.rgb = COLOR_KPI_VALUE
        run_v.font.name = FONT_NAME_JP

        # 補足 (下部)
        if sub:
            sub_box = slide.shapes.add_textbox(
                cx + Inches(0.15), cy + card_h - Inches(0.40),
                int(card_w - Inches(0.30)), Inches(0.28)
            )
            tf_s = sub_box.text_frame
            tf_s.word_wrap = True
            p_s = tf_s.paragraphs[0]
            p_s.alignment = PP_ALIGN.CENTER
            run_s = p_s.add_run()
            run_s.text = sub
            run_s.font.size = Pt(KPI_SUB_PT)
            run_s.font.color.rgb = COLOR_KPI_SUB
            run_s.font.name = FONT_NAME_JP

    print(f"  ✓ KPI カード: {n}枚 ({n_cols}列×{n_rows}行)")


def add_source_label(slide, text):
    """出典: roleup なら Source 3 placeholder、stella なら動的 textbox。"""
    src_shape = find_shape(slide, SHAPE_SOURCE, warn=False)
    if src_shape is not None:
        tf = src_shape.text_frame
        tf.word_wrap = True
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)
        p = tf.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(SOURCE_FONT_PT)
        run.font.color.rgb = COLOR_SOURCE
        run.font.name = FONT_NAME_JP
        print(f"  ✓ 出典 ({SHAPE_SOURCE} placeholder): {text[:50]}")
        return

    txBox = slide.shapes.add_textbox(SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(SOURCE_FONT_PT)
    run.font.color.rgb = COLOR_SOURCE
    run.font.name = FONT_NAME_JP
    print(f"  ✓ 出典 (textbox): {text[:50]}")


def main():
    parser = argparse.ArgumentParser(description="事業セグメント概要 PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True)
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    parser.add_argument("--output", required=True)
    add_brand_arg(parser)
    args = parser.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "business-overview")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    print(f"=== 事業セグメント概要スライド生成 (brand={theme.id}) ===")
    print(f"  Template: {template_path}")

    _validate(data)
    require_source(data, theme, skill_id=SKILL_ID)

    # chart_title のデフォルト値を data に埋めて、brand 別の top/subtitle 解決で
    # どちらの placeholder field に書かれても落ちない状態にする。
    parent_company = data["parent_company"]
    segment_name = data["segment_name"]
    if not data.get("chart_title"):
        data["chart_title"] = f"{parent_company}:{segment_name}の概要"

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top placeholder (brand-aware)
    top_text = resolve_top_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:50]}")

    # Subtitle placeholder
    sub_text = resolve_subtitle_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:50]}")

    # Stella V1: Table 1 placeholder removal (silent for roleup which has no Table 1)
    _silent_remove_shape(slide, "Table 1")

    # Roleup: silently remove brown guide rectangles
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # Left panel: 事業の概要
    overview = data["overview"]
    section_title_left = overview.get("section_title", "事業の概要")
    add_section_title(slide, section_title_left, LEFT_X, PANEL_Y, LEFT_W)
    build_overview_table(slide, overview["items"], LEFT_X, PANEL_Y + Inches(0.40), LEFT_W)

    # Right panel: 業績 or 主要 KPI
    perf = data["performance"]
    mode = perf["mode"]
    section_title_right = perf.get(
        "section_title",
        "業績" if mode == "revenue_chart" else "主要 KPI",
    )
    add_section_title(slide, section_title_right, RIGHT_X, PANEL_Y, RIGHT_W)

    if mode == "revenue_chart":
        unit_label = perf.get("unit_label", "")
        legend_max_right = None
        unit_w = Inches(1.50)
        if unit_label:
            # 単位ラベルは右上 (RIGHT_X+RIGHT_W-unit_w から開始)、凡例は左寄りに配置
            unit_label_left = RIGHT_X + RIGHT_W - unit_w
            add_unit_label(slide, unit_label, unit_label_left, PANEL_Y + Inches(0.35),
                           unit_w, height=Inches(0.27))
            legend_max_right = unit_label_left - Inches(0.10)

        add_custom_legend(slide, perf, RIGHT_X, PANEL_Y + Inches(0.35), RIGHT_W,
                          max_right_emu=legend_max_right)

        chart_top = PANEL_Y + Inches(0.55)
        build_combo_chart(slide, perf, RIGHT_X, chart_top, RIGHT_W, CHART_H)
        add_cagr_annotation(slide, perf, RIGHT_X, chart_top, RIGHT_W, CHART_H)
    else:  # kpi_cards
        cards_top = PANEL_Y + Inches(0.55)
        build_kpi_cards(slide, perf, RIGHT_X, cards_top, RIGHT_W, CHART_H)

    # Source
    source = data.get("source", "")
    if source:
        body = source if source.startswith("出典") else f"出典:{source}"
        add_source_label(slide, body)

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
