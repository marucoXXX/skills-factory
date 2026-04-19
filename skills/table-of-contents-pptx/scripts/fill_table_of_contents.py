"""
fill_table_of_contents.py — 目次（Table of Contents）スライドを生成

レイアウト:
  - 上部: 「目次」+ サブタイトル
  - 中央: セクション番号 + セクション名 + (オプション) サブ項目リスト
      各セクションは色付きバッジ + Boldタイトル + 薄色サブテキスト

Usage:
  python fill_table_of_contents.py \
    --data /home/claude/toc_data.json \
    --template <path>/table-of-contents-pptx-template.pptx \
    --output /mnt/user-data/outputs/TableOfContents_output.pptx
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

TOC_LEFT = Inches(0.80)
TOC_TOP = Inches(1.55)
TOC_WIDTH = Inches(11.73)
TOC_HEIGHT = Inches(5.40)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUBTEXT = RGBColor(0x66, 0x66, 0x66)
COLOR_DIVIDER = RGBColor(0xDD, 0xDD, 0xDD)

# セクション色のローテーション（PEST色と統一感）
SECTION_COLORS = [
    RGBColor(0x2E, 0x4A, 0x6B),   # 紺
    RGBColor(0x7B, 0x4F, 0xB0),   # 紫
    RGBColor(0x2E, 0x6F, 0xBF),   # 青
    RGBColor(0x3D, 0x8F, 0x5A),   # 緑
    RGBColor(0xDA, 0x7A, 0x2D),   # オレンジ
    RGBColor(0xC0, 0x3A, 0x3A),   # 赤
    RGBColor(0x59, 0x59, 0x59),   # グレー
]

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_NUMBER = Pt(28)
FONT_SIZE_SECTION = Pt(18)
FONT_SIZE_SUBITEM = Pt(11)
FONT_SIZE_PAGE = Pt(11)
FONT_SIZE_SOURCE = Pt(10)


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name=FONT_NAME_JP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def draw_section_row(slide, idx, section, left, top, width, height):
    """1セクション行を描画"""
    # 色決定
    color_hex = section.get("color")
    if color_hex:
        color = hex_to_rgb(color_hex)
    else:
        color = SECTION_COLORS[(idx - 1) % len(SECTION_COLORS)]

    # 番号バッジ（左、四角）
    badge_w = Inches(0.85)
    badge_h = height - Inches(0.10)
    badge_y = top + Inches(0.05)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, badge_y, badge_w, badge_h,
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.shadow.inherit = False

    # 番号テキスト
    tf = badge.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    p_elem = etree.SubElement(tf._txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "ctr")
    r = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r, qn("a:rPr"), attrib={
        "lang": "en-US",
        "sz": str(int(FONT_SIZE_NUMBER.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": "Arial"})
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    s = etree.SubElement(sf, qn("a:srgbClr"))
    s.set("val", "FFFFFF")
    t = etree.SubElement(r, qn("a:t"))
    t.text = f"{idx:02d}"

    # コンテンツ領域
    content_left = left + badge_w + Inches(0.20)
    content_w = width - badge_w - Inches(0.20) - Inches(0.80)  # 右にページ番号領域
    page_x = left + width - Inches(0.80)
    page_w = Inches(0.80)

    # セクションタイトル
    title = section.get("title", f"セクション{idx}")
    subitems = section.get("subitems", [])
    has_subitems = len(subitems) > 0

    if has_subitems:
        # タイトルを上半分、サブ項目を下半分
        title_y = top + Inches(0.05)
        title_h = Inches(0.40)
    else:
        # タイトルを中央寄せ
        title_y = top
        title_h = height

    # タイトル + サブ項目を1つのテキストボックスで管理
    tb = slide.shapes.add_textbox(content_left, top + Inches(0.05), content_w, height - Inches(0.10))
    tf2 = tb.text_frame
    tf2.word_wrap = True
    tf2.margin_left = 0; tf2.margin_right = 0
    tf2.margin_top = 0; tf2.margin_bottom = 0
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE if not has_subitems else MSO_ANCHOR.TOP

    for p in list(tf2.paragraphs):
        p._p.getparent().remove(p._p)

    # タイトル段落
    p_title = etree.SubElement(tf2._txBody, qn("a:p"))
    pPr_t = etree.SubElement(p_title, qn("a:pPr"))
    pPr_t.set("algn", "l")
    r_title = etree.SubElement(p_title, qn("a:r"))
    rPr_title = etree.SubElement(r_title, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(FONT_SIZE_SECTION.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr_title, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr_title, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    sf_title = etree.SubElement(rPr_title, qn("a:solidFill"))
    s_title = etree.SubElement(sf_title, qn("a:srgbClr"))
    s_title.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))
    t_title = etree.SubElement(r_title, qn("a:t"))
    t_title.text = title

    # サブ項目段落
    if has_subitems:
        for sub in subitems:
            p_sub = etree.SubElement(tf2._txBody, qn("a:p"))
            pPr_sub = etree.SubElement(p_sub, qn("a:pPr"), attrib={
                "marL": "180000",
                "indent": "-180000",
            })
            spcBef = etree.SubElement(pPr_sub, qn("a:spcBef"))
            etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "200"})

            buChar = etree.SubElement(pPr_sub, qn("a:buChar"), attrib={"char": "▸"})
            buFont = etree.SubElement(pPr_sub, qn("a:buFont"), attrib={"typeface": "Arial"})
            buClr = etree.SubElement(pPr_sub, qn("a:buClr"))
            buClrSolid = etree.SubElement(buClr, qn("a:srgbClr"))
            buClrSolid.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))

            r_sub = etree.SubElement(p_sub, qn("a:r"))
            rPr_sub = etree.SubElement(r_sub, qn("a:rPr"), attrib={
                "lang": "ja-JP",
                "sz": str(int(FONT_SIZE_SUBITEM.pt * 100)),
                "b": "0",
            })
            etree.SubElement(rPr_sub, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
            etree.SubElement(rPr_sub, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
            sf_sub = etree.SubElement(rPr_sub, qn("a:solidFill"))
            s_sub = etree.SubElement(sf_sub, qn("a:srgbClr"))
            s_sub.set("val", "555555")
            t_sub = etree.SubElement(r_sub, qn("a:t"))
            t_sub.text = sub

    # ページ番号（右端）
    page = section.get("page", "")
    if page:
        add_text_box(
            slide, f"P. {page}",
            page_x, top, page_w, height,
            FONT_SIZE_PAGE, bold=True,
            color=color,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )

    # 区切り線（行の下）
    div_y = top + height - Emu(int(Inches(0.01)))
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + badge_w + Inches(0.20), div_y,
        width - badge_w - Inches(0.20), Emu(int(Inches(0.01))),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_DIVIDER
    div.line.fill.background()


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument("--template", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation(args.template)
    slide = prs.slides[0]

    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), data.get("main_message", "目次"))
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), data.get("chart_title", "Table of Contents"))
    print(f"  ✓ Title set")

    sections = data.get("sections", [])
    if not sections:
        print("  ✗ ERROR: 'sections' is required", file=sys.stderr)
        sys.exit(1)

    n = len(sections)
    if n > 7:
        print(f"  ⚠ WARNING: {n} sections > 7. Only first 7 will be shown.", file=sys.stderr)
        sections = sections[:7]
        n = 7

    # 行の高さを計算（セクション数に応じて）
    gap = Inches(0.10)
    section_h = Emu(int((TOC_HEIGHT - gap * (n - 1)) / n))

    for i, section in enumerate(sections):
        y = TOC_TOP + (section_h + gap) * i
        draw_section_row(slide, i + 1, section, TOC_LEFT, y, TOC_WIDTH, section_h)
        print(f"  ✓ Section {i+1}: {section.get('title', '')[:40]}")

    source = data.get("source", "")
    if source:
        add_text_box(
            slide, source,
            SOURCE_X, SOURCE_Y, SOURCE_W, Inches(0.25),
            FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
            align=PP_ALIGN.LEFT,
        )

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
