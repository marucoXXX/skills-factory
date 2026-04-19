"""
fill_executive_summary.py — エグゼクティブサマリースライドをPPTXネイティブオブジェクトで生成

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 中央: 3〜5個のKey Findingsを縦積みで表示
      各Finding: 番号バッジ + カテゴリラベル + 見出し + 詳細テキスト
  - 下部: 出典

Usage:
  python fill_executive_summary.py \
    --data /home/claude/executive_summary_data.json \
    --template <path>/executive-summary-template.pptx \
    --output /mnt/user-data/outputs/ExecutiveSummary_output.pptx
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ── Layout Constants ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

# Findingsグリッド配置
GRID_LEFT = Inches(0.41)
GRID_TOP = Inches(1.55)
GRID_WIDTH = Inches(12.51)
GRID_HEIGHT = Inches(5.35)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(6.93)
SOURCE_W = Inches(12.50)

# ── Colors ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUBTEXT = RGBColor(0x55, 0x55, 0x55)

# カテゴリ色マッピング（findingの category フィールドで使用）
CATEGORY_COLORS = {
    "対象会社": RGBColor(0x2E, 0x4A, 0x6B),       # 紺
    "company": RGBColor(0x2E, 0x4A, 0x6B),
    "target": RGBColor(0x2E, 0x4A, 0x6B),
    "マクロ環境": RGBColor(0x7B, 0x4F, 0xB0),      # 紫
    "macro": RGBColor(0x7B, 0x4F, 0xB0),
    "pest": RGBColor(0x7B, 0x4F, 0xB0),
    "市場": RGBColor(0x2E, 0x6F, 0xBF),           # 青
    "market": RGBColor(0x2E, 0x6F, 0xBF),
    "競合": RGBColor(0xDA, 0x7A, 0x2D),           # オレンジ
    "competitor": RGBColor(0xDA, 0x7A, 0x2D),
    "財務": RGBColor(0x3D, 0x8F, 0x5A),           # 緑
    "financial": RGBColor(0x3D, 0x8F, 0x5A),
    "リスク": RGBColor(0xB8, 0x3A, 0x3A),         # 赤
    "risk": RGBColor(0xB8, 0x3A, 0x3A),
    "機会": RGBColor(0x1B, 0x7A, 0x3B),           # 濃緑
    "opportunity": RGBColor(0x1B, 0x7A, 0x3B),
    "示唆": RGBColor(0x55, 0x55, 0x55),           # グレー
    "implication": RGBColor(0x55, 0x55, 0x55),
    "結論": RGBColor(0x33, 0x33, 0x33),           # 濃グレー
    "conclusion": RGBColor(0x33, 0x33, 0x33),
}

DEFAULT_COLOR = RGBColor(0x4E, 0x79, 0xA7)  # デフォルト: 紺系

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_BADGE = Pt(20)
FONT_SIZE_CATEGORY = Pt(10)
FONT_SIZE_HEADING = Pt(14)
FONT_SIZE_DETAIL = Pt(11)
FONT_SIZE_SOURCE = Pt(10)


# ──────────────────────────────────────────────
# Utility
# ──────────────────────────────────────────────
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 italic=False, font_name=FONT_NAME_JP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_category_color(category):
    """カテゴリ名から色を取得。未知のカテゴリはデフォルト色"""
    if not category:
        return DEFAULT_COLOR
    # 完全一致
    if category in CATEGORY_COLORS:
        return CATEGORY_COLORS[category]
    # 小文字でも試す
    lower = category.lower()
    if lower in CATEGORY_COLORS:
        return CATEGORY_COLORS[lower]
    return DEFAULT_COLOR


# ──────────────────────────────────────────────
# Number Badge (circular)
# ──────────────────────────────────────────────
def draw_number_badge(slide, number, color, left, top, diameter):
    """番号入りの円形バッジを描画"""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, diameter, diameter,
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.shadow.inherit = False

    # テキスト設定（番号）
    tf = badge.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 既存段落クリア
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    p_elem = etree.SubElement(tf._txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "ctr")

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "en-US",
        "sz": str(int(FONT_SIZE_BADGE.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": "Arial"})
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    s = etree.SubElement(sf, qn("a:srgbClr"))
    s.set("val", "FFFFFF")
    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = f"{number:02d}"


# ──────────────────────────────────────────────
# Finding Row
# ──────────────────────────────────────────────
def draw_finding(slide, idx, finding, left, top, width, height):
    """
    1つのFindingを描画する
    レイアウト:
      [01] [カテゴリラベル] [見出し (Bold)]
           [詳細テキスト                ]
    """
    # 番号バッジ
    badge_size = Inches(0.55)
    badge_left = left
    badge_top = top + Inches(0.05)

    category = finding.get("category", "")
    color = None
    color_hex = finding.get("color")
    if color_hex:
        color = hex_to_rgb(color_hex)
    else:
        color = get_category_color(category)

    draw_number_badge(slide, idx, color, badge_left, badge_top, badge_size)

    # 左側の縦バー（カテゴリ色）
    bar_left = left + badge_size + Inches(0.10)
    bar_top = top
    bar_w = Inches(0.05)
    bar_h = height - Inches(0.10)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_w, bar_h,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    bar.text_frame.text = ""

    # テキスト領域開始
    content_left = bar_left + bar_w + Inches(0.15)
    content_w = width - (content_left - left)

    # 上段: カテゴリラベル + 見出し
    top_row_top = top + Inches(0.02)
    top_row_h = Inches(0.35)

    # カテゴリラベル（小さな色付きタグ風）
    if category:
        # タグ幅を動的に決定（固定幅でOK）
        cat_label_w = Inches(1.30)
        cat_tb = slide.shapes.add_textbox(
            content_left, top_row_top, cat_label_w, top_row_h,
        )
        ctf = cat_tb.text_frame
        ctf.margin_left = 0; ctf.margin_right = 0
        ctf.margin_top = 0; ctf.margin_bottom = 0
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

        for p in list(ctf.paragraphs):
            p._p.getparent().remove(p._p)

        p_elem = etree.SubElement(ctf._txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        pPr.set("algn", "l")

        r_elem = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(FONT_SIZE_CATEGORY.pt * 100)),
            "b": "1",
        })
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        s = etree.SubElement(sf, qn("a:srgbClr"))
        s.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = f"▍ {category}"

        heading_left = content_left + cat_label_w
        heading_w = content_w - cat_label_w
    else:
        heading_left = content_left
        heading_w = content_w

    # 見出し
    heading = finding.get("heading", "")
    if heading:
        add_text_box(
            slide, heading,
            heading_left, top_row_top, heading_w, top_row_h,
            FONT_SIZE_HEADING, bold=True,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )

    # 下段: 詳細テキスト
    detail = finding.get("detail", "")
    if detail:
        detail_top = top_row_top + top_row_h + Inches(0.05)
        detail_h = height - (detail_top - top) - Inches(0.10)
        add_text_box(
            slide, detail,
            content_left, detail_top, content_w, detail_h,
            FONT_SIZE_DETAIL, bold=False,
            color=COLOR_SUBTEXT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument("--template", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation(args.template)
    slide = prs.slides[0]

    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), data.get("main_message", ""))
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), data.get("chart_title", "エグゼクティブサマリー"))
    print(f"  ✓ Main Message & Chart Title set")

    findings = data.get("findings", [])
    if not findings:
        print("  ✗ ERROR: 'findings' is required", file=sys.stderr)
        sys.exit(1)

    n = len(findings)
    if n > 6:
        print(f"  ⚠ WARNING: {n} findings > 6. Only first 6 will be shown.", file=sys.stderr)
        findings = findings[:6]
        n = 6

    # 各Findingの高さを動的に計算
    # ギャップは等間隔
    total_h = GRID_HEIGHT
    gap_h = Inches(0.12)
    total_gap = gap_h * (n - 1) if n > 1 else Emu(0)
    finding_h = Emu(int((total_h - total_gap) / n))

    for i, f in enumerate(findings):
        top = GRID_TOP + (finding_h + gap_h) * i
        draw_finding(slide, i + 1, f, GRID_LEFT, top, GRID_WIDTH, finding_h)
        print(f"  ✓ Finding {i+1}: {f.get('heading', '')[:40]}...")

    # 出典
    source = data.get("source", "")
    if source:
        add_text_box(
            slide, source,
            SOURCE_X, SOURCE_Y, SOURCE_W, Inches(0.25),
            FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
            align=PP_ALIGN.LEFT,
        )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
