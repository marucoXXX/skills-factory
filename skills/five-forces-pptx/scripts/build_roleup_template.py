"""One-shot generator: five-forces-pptx の roleup curated template を
stella template から派生して作る。

実行は本ファイルを直接 python3 で起動するだけ。生成済テンプレは
`assets/roleup/five-forces-template.pptx` に上書き保存する。

設計:
  - スライドサイズ: A4 landscape (10691813 × 7559675 EMU = 11.69 × 8.27 in)
  - Title 1 / Text Placeholder 2: roleup pilot (cp / business-model) と同じ座標へ
  - Rectangle 5/15/17/18/20 + Straight Arrow Connector 22/23/27/28: 比率スケール (W: ×11.69/13.33, H: ×8.27/7.50)
  - TextBox 30 (意味合い): 右端配置を維持しつつスケール
  - Source 3: roleup 慣習通り下端に新規追加 (TextBox)
  - フォント: Yu Gothic UI (Title / Subtitle / 各 Rectangle 内テキスト / 意味合い / Source)
  - think-cell OLE は削除
"""
import copy
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(SKILL_DIR, "assets", "stellar_aiz", "five-forces-template.pptx")
DST = os.path.join(SKILL_DIR, "assets", "roleup", "five-forces-template.pptx")

# A4 landscape EMU
A4_W = 10691813
A4_H = 7559675

# stella source size in EMU
STELLA_W = Emu(int(13.33 * 914400))
STELLA_H = Emu(int(7.50 * 914400))

FONT_EA = "Yu Gothic UI"
COLOR_TEXT = RGBColor(0x24, 0x1A, 0x17)
COLOR_SUBTITLE = RGBColor(0x89, 0x71, 0x41)
COLOR_SOURCE = RGBColor(0x3E, 0x3A, 0x39)


def _find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def _set_run_font(run, font_name, color, size_pt=None, bold=False):
    run.font.name = font_name
    run.font.color.rgb = color
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    run.font.bold = bold
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn("a:rPr"))
    for tag in ("a:latin", "a:ea"):
        el = rPr.find(qn(tag))
        if el is not None:
            rPr.remove(el)
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": font_name})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": font_name})


def _retypeset_text(shape, color, size_pt=None, bold=False):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            _set_run_font(run, FONT_EA, color, size_pt=size_pt, bold=bold)


def main():
    prs = Presentation(SRC)
    prs.slide_width = A4_W
    prs.slide_height = A4_H

    slide = prs.slides[0]

    scale_w = A4_W / STELLA_W
    scale_h = A4_H / STELLA_H

    # 1. Remove think-cell OLE
    for s in list(slide.shapes):
        if "think-cell" in (s.name or ""):
            sp = s._element
            sp.getparent().remove(sp)

    # 2. Title 1 (placeholder) — 固定座標 (cp/business-model 流儀)
    title = _find(slide, "Title 1")
    if title is not None:
        title.left = Inches(0.41)
        title.top = Inches(0.41)
        title.width = Inches(10.87)
        title.height = Inches(0.50)
        _retypeset_text(title, COLOR_TEXT, size_pt=22, bold=True)

    # 3. Text Placeholder 2 (subtitle) — 固定座標
    sub = _find(slide, "Text Placeholder 2")
    if sub is not None:
        sub.left = Inches(0.41)
        sub.top = Inches(1.00)
        sub.width = Inches(10.87)
        sub.height = Inches(0.36)
        _retypeset_text(sub, COLOR_SUBTITLE, size_pt=12, bold=False)

    # 4. Rectangle 5/15/17/18/20 + 矢印 4 本: 比率スケール
    scaled_targets = [
        "Rectangle 5", "Rectangle 15", "Rectangle 17", "Rectangle 18", "Rectangle 20",
        "Straight Arrow Connector 22", "Straight Arrow Connector 23",
        "Straight Arrow Connector 27", "Straight Arrow Connector 28",
    ]
    for name in scaled_targets:
        sh = _find(slide, name)
        if sh is None:
            continue
        sh.left = int(sh.left * scale_w)
        sh.top = int(sh.top * scale_h)
        sh.width = int(sh.width * scale_w)
        sh.height = int(sh.height * scale_h)
        if sh.has_text_frame:
            _retypeset_text(sh, COLOR_TEXT, size_pt=10, bold=False)

    # 5. TextBox 30 (意味合い): 右端パネル
    tb30 = _find(slide, "TextBox 30")
    if tb30 is not None:
        tb30.left = Inches(8.00)
        tb30.top = Inches(1.55)
        tb30.width = Inches(3.28)
        tb30.height = Inches(5.60)
        if tb30.has_text_frame:
            for pi, para in enumerate(tb30.text_frame.paragraphs):
                for run in para.runs:
                    if pi == 0:
                        _set_run_font(run, FONT_EA, COLOR_TEXT, size_pt=12, bold=True)
                    else:
                        _set_run_font(run, FONT_EA, COLOR_TEXT, size_pt=10, bold=False)

    # 6. Source 3 — 下端に新規追加
    if _find(slide, "Source 3") is None:
        src_tb = slide.shapes.add_textbox(
            Inches(0.41), Inches(7.45), Inches(10.87), Inches(0.40)
        )
        src_tb.name = "Source 3"
        tf = src_tb.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "出典: "
        _set_run_font(run, FONT_EA, COLOR_SOURCE, size_pt=6, bold=False)

    prs.save(DST)
    print(f"Generated: {DST}")
    print(f"  slide_size: {prs.slide_width} × {prs.slide_height} EMU")
    for s in slide.shapes:
        try:
            l, t, w, h = s.left/914400, s.top/914400, s.width/914400, s.height/914400
        except Exception:
            l = t = w = h = 0
        print(f"  - {s.shape_type} \"{s.name}\" L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")


if __name__ == "__main__":
    main()
