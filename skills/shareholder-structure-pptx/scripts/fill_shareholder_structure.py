"""
fill_shareholder_structure.py — 株主構成・役員構成をPPTXネイティブテーブルとして生成

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。

テンプレート構造:
  Stella (assets/stellar_aiz/shareholder-structure-template.pptx, 16:9):
    - Title 1            (PLACEHOLDER): メインメッセージ
    - Text Placeholder 2 (PLACEHOLDER): チャートタイトル
    - Table 1            (TABLE):       テンプレート用テーブル → 削除
  Roleup (assets/roleup/shareholder-structure-template.pptx, A4 横, ch ベース):
    - Title 1, Text Placeholder 2, Source 3 (PLACEHOLDER)
    - 茶色ガイド `正方形/長方形 1` `正方形/長方形 8` (object 8) は fill が silent_remove

生成物:
  - ■株主構成 セクションタイトル (TextBox)
  - 株主テーブル (ネイティブTable: 7列)
  - ■役員構成 セクションタイトル (TextBox)
  - 役員テーブル (ネイティブTable: 6列)
  - 出典 (Source 3 placeholder for roleup, dynamic TextBox fallback for stella)

Usage:
  python fill_shareholder_structure.py --brand stellar_aiz \\
    --data {{WORK_DIR}}/shareholder_structure_data.json \\
    --output {{OUTPUT_DIR}}/ShareholderStructure_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "shareholder-structure-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_TABLE = "Table 1"

# Defaults (stella). Reassigned in _apply_theme(theme) for roleup.
SHAPE_SOURCE = "Source"            # roleup uses 'Source 3'
TABLE_LEFT = Inches(0.41)
TABLE_WIDTH = Inches(12.52)
CONTENT_START_Y = Inches(1.30)
BOTTOM_MARGIN = Inches(0.30)

FONT_NAME_JP = "Meiryo UI"
TEXT_HEX = "333333"
HEADER_BG_HEX = "F5F0D0"           # stella: テンプレ準拠ベージュ
TOTAL_BG_HEX = "F0F0F0"            # stella: 合計行
BORDER_HEX = "CCCCCC"              # stella: 罫線
SECTION_TITLE_PT = 14
SECTION_TITLE_BOLD = True
SECTION_TITLE_HEX = "333333"
SOURCE_PT = 10
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.10)
SOURCE_W = Inches(10.0)
SOURCE_H = Inches(0.25)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)

# Font candidates (hundredths of a pt). Reassigned per brand: roleup must
# stay in C4 allowed set {22, 14, 12, 10, 6} so candidates collapse to (1000,).
FONT_CANDIDATES = (1300, 1200, 1100, 1000, 900)
MIN_ROW_HEIGHT = Inches(0.25)
HEADER_HEIGHT = Inches(0.35)
MAX_ROW_HEIGHT = Inches(0.40)

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global SHAPE_SOURCE
    global TABLE_LEFT, TABLE_WIDTH, CONTENT_START_Y, BOTTOM_MARGIN
    global FONT_NAME_JP, TEXT_HEX, HEADER_BG_HEX, TOTAL_BG_HEX, BORDER_HEX
    global SECTION_TITLE_PT, SECTION_TITLE_BOLD, SECTION_TITLE_HEX
    global SOURCE_PT, SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE, FONT_CANDIDATES

    _THEME = theme

    FONT_NAME_JP = theme.font_ea
    TEXT_HEX = theme.hex_no_hash("text")
    HEADER_BG_HEX = theme.hex_no_hash("header_bg")
    BORDER_HEX = theme.hex_no_hash("highlight_other")

    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        TABLE_LEFT = Inches(0.41)
        TABLE_WIDTH = Inches(12.52)
        CONTENT_START_Y = Inches(1.30)
        BOTTOM_MARGIN = Inches(0.30)
        TOTAL_BG_HEX = "F0F0F0"
        SECTION_TITLE_PT = 14
        SECTION_TITLE_BOLD = True
        SECTION_TITLE_HEX = TEXT_HEX
        SOURCE_PT = 10
        SOURCE_X = Inches(theme.layout_in("source_x_in")) if "source_x_in" in theme._layout else Inches(0.41)
        SOURCE_Y = Inches(theme.layout_in("source_y_in")) if "source_y_in" in theme._layout else Inches(7.10)
        SOURCE_W = Inches(theme.layout_in("source_w_in")) if "source_w_in" in theme._layout else Inches(10.0)
        SOURCE_H = Inches(theme.layout_in("source_h_in")) if "source_h_in" in theme._layout else Inches(0.25)
        FONT_CANDIDATES = (1300, 1200, 1100, 1000, 900)
    else:
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        SHAPE_SOURCE = "Source 3"
        # A4 横テンプレ — left guide x=0.41, width = 11.693 - 2*0.41 = 10.873 in
        TABLE_LEFT = Inches(theme.layout_rule("left_align_guide_x_in", 0.41))
        slide_w_in = theme.slide_w / 914400
        TABLE_WIDTH = Inches(slide_w_in - 2 * theme.layout_rule("left_align_guide_x_in", 0.41))
        CONTENT_START_Y = Inches(1.40)
        BOTTOM_MARGIN = Inches(0.45)  # Source 3 placeholder のためのマージン
        TOTAL_BG_HEX = theme.hex_no_hash("label_bg").upper()
        SECTION_TITLE_PT = theme.pt_value("font_size_subtitle_pt")  # 12pt
        SECTION_TITLE_BOLD = False
        SECTION_TITLE_HEX = theme.hex_no_hash("subtitle")  # #897141
        SOURCE_PT = theme.pt_value("font_size_source_pt")  # 6pt
        SOURCE_X = theme.layout("source_x_in")
        SOURCE_Y = theme.layout("source_y_in")
        SOURCE_W = theme.layout("source_w_in")
        SOURCE_H = theme.layout("source_h_in")
        # All body/header text must be 10pt for roleup (C4)
        FONT_CANDIDATES = (1000,)


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_shape(slide, name, warn=True):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    if warn:
        print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    """TextBox のテキストを上書きし、フォントを theme.font_ea に強制設定。"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs:
            rPr = run._r.find(qn("a:rPr"))
            if rPr is None:
                rPr = etree.SubElement(run._r, qn("a:rPr"), attrib={"lang": "ja-JP"})
                run._r.insert(0, rPr)
            for tag in [qn("a:latin"), qn("a:ea")]:
                old = rPr.find(tag)
                if old is not None:
                    rPr.remove(old)
                etree.SubElement(rPr, tag, attrib={"typeface": FONT_NAME_JP})
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_section_title(slide, text, left, top, width):
    """セクションタイトル (■株主構成 等) を TextBox で追加。"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(SECTION_TITLE_PT)
    run.font.bold = SECTION_TITLE_BOLD
    # color: stella=text, roleup=subtitle (#897141)
    r, g, b = int(SECTION_TITLE_HEX[0:2], 16), int(SECTION_TITLE_HEX[2:4], 16), int(SECTION_TITLE_HEX[4:6], 16)
    run.font.color.rgb = RGBColor(r, g, b)
    run.font.name = FONT_NAME_JP
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        old_ea = rPr.find(qn("a:ea"))
        if old_ea is not None:
            rPr.remove(old_ea)
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    print(f"  ✓ セクションタイトル: {text}")
    return txBox


def add_dynamic_source_textbox(slide, text):
    """Fallback: Source 3 placeholder が無い場合の動的 textbox 配置 (stella)。"""
    txBox = slide.shapes.add_textbox(SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(SOURCE_PT)
    run.font.color.rgb = COLOR_SOURCE
    run.font.name = FONT_NAME_JP
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        old_ea = rPr.find(qn("a:ea"))
        if old_ea is not None:
            rPr.remove(old_ea)
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    print(f"  ✓ 出典 (dynamic textbox): {text}")


def write_source_placeholder(shape, text):
    """Source 3 placeholder にテキストを書き込み、フォント size/name を強制 (roleup)。"""
    tf = shape.text_frame
    para = tf.paragraphs[0]
    for r in list(para.runs):
        r.text = ""
    if para.runs:
        run = para.runs[0]
        run.text = text
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text
        run = para.runs[0]
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        run._r.insert(0, rPr)
    rPr.set("sz", str(SOURCE_PT * 100))
    for tag in [qn("a:latin"), qn("a:ea")]:
        old = rPr.find(tag)
        if old is not None:
            rPr.remove(old)
        etree.SubElement(rPr, tag, attrib={"typeface": FONT_NAME_JP})


def apply_cell(cell, text, is_header=False, bold=False, align="l",
               font_size=1200, bg_hex=None):
    """セルにテキスト・スタイルを設定 (ネイティブPPTXセル)。"""
    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        txBody = etree.SubElement(tc, qn("a:txBody"))

    old_bodyPr = txBody.find(qn("a:bodyPr"))
    if old_bodyPr is not None:
        txBody.remove(old_bodyPr)
    bodyPr = etree.SubElement(txBody, qn("a:bodyPr"), attrib={
        "wrap": "square",
        "lIns": "0", "rIns": "0",
        "tIns": "0", "bIns": "0",
        "anchor": "ctr",
    })
    txBody.insert(0, bodyPr)

    if txBody.find(qn("a:lstStyle")) is None:
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        txBody.insert(1, lstStyle)

    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    p_elem = etree.SubElement(txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", align)

    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPct"), attrib={"val": "100000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "0"})
    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    etree.SubElement(spcAft, qn("a:spcPts"), attrib={"val": "0"})

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(font_size),
        "b": "1" if (bold or is_header) else "0",
    })
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr"), attrib={"val": TEXT_HEX})
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = str(text)

    old_tcPr = tc.find(qn("a:tcPr"))
    if old_tcPr is not None:
        tc.remove(old_tcPr)

    tcPr = etree.SubElement(tc, qn("a:tcPr"), attrib={
        "marL": "36576", "marR": "36576",
        "marT": "18288", "marB": "18288",
        "anchor": "ctr",
    })

    for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = etree.SubElement(tcPr, qn(border_name), attrib={"w": "6350", "cmpd": "sng"})
        sf = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": BORDER_HEX})

    if bg_hex:
        cell_fill = etree.SubElement(tcPr, qn("a:solidFill"))
        etree.SubElement(cell_fill, qn("a:srgbClr"), attrib={"val": bg_hex})


def build_table(slide, columns, rows, col_widths_inch, left, top, width,
                font_size=1200, has_total=False, row_height=Inches(0.30)):
    """ネイティブPPTXテーブルを構築する。"""
    n_rows = len(rows) + 1
    n_cols = len(columns)

    table_height = HEADER_HEIGHT + row_height * len(rows)

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = shape.table

    for i, w in enumerate(col_widths_inch):
        table.columns[i].width = Inches(w)

    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    for i, tr in enumerate(tbl_elem.findall(qn('a:tr'))):
        if i == 0:
            tr.set('h', str(HEADER_HEIGHT))
        else:
            tr.set('h', str(row_height))

    print(f"  ✓ テーブル高さ: {table_height/914400:.2f}in "
          f"(データ行: {row_height/914400:.2f}in × {len(rows)}行)")

    for c_idx, col_name in enumerate(columns):
        align = "r" if col_name in ["持株数", "議決権比率(%)", "役員報酬"] else "ctr" if col_name == "#" else "l"
        apply_cell(table.cell(0, c_idx), col_name, is_header=True,
                   align=align, font_size=font_size, bg_hex=HEADER_BG_HEX)

    for r_idx, row_data in enumerate(rows):
        cells = row_data["cells"]
        aligns = row_data.get("aligns", ["l"] * n_cols)
        is_total = has_total and (r_idx == len(rows) - 1)
        bg = TOTAL_BG_HEX if is_total else None

        for c_idx, cell_text in enumerate(cells):
            align = aligns[c_idx] if c_idx < len(aligns) else "l"
            apply_cell(table.cell(r_idx + 1, c_idx), cell_text,
                       bold=is_total, align=align, font_size=font_size,
                       bg_hex=bg)

    print(f"  ✓ テーブル生成: {n_rows}行 × {n_cols}列")
    return shape, table_height


def format_shareholder_rows(data):
    rows = []
    aligns = ["ctr", "l", "l", "l", "r", "r", "l"]

    for row in data.get("rows", []):
        cells = [
            str(row.get("number", "")),
            row.get("name", ""),
            row.get("position", ""),
            row.get("relation", ""),
            row.get("shares", ""),
            f'{row.get("voting_ratio", 0):.1f}',
            row.get("note", ""),
        ]
        rows.append({"cells": cells, "aligns": aligns})

    total = data.get("total")
    if total:
        cells = [
            "", "合計", "", "",
            total.get("shares", ""),
            f'{total.get("voting_ratio", 100.0):.1f}',
            "",
        ]
        rows.append({"cells": cells, "aligns": aligns})

    return rows


def format_director_rows(data):
    rows = []
    aligns = ["ctr", "l", "l", "l", "r", "l"]

    for row in data.get("rows", []):
        cells = [
            str(row.get("number", "")),
            row.get("name", ""),
            row.get("position", ""),
            row.get("relation", ""),
            row.get("compensation", ""),
            row.get("note", ""),
        ]
        rows.append({"cells": cells, "aligns": aligns})

    return rows


def required_row_height(font_size_hundredths):
    """フォントサイズから、レンダリング後の実際の行高さ(Emu)を算出する。"""
    font_pt = font_size_hundredths / 100
    height_in = (font_pt * 1.5) / 72 + 0.04 + 0.05
    return Inches(height_in)


def main():
    parser = argparse.ArgumentParser(description="株主構成・役員構成 PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True, help="JSONデータファイルパス")
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    parser.add_argument("--output", required=True, help="出力PPTXファイルパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "shareholder-structure")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    print(f"=== 株主構成・役員構成スライド生成 (brand={theme.id}) ===")
    print(f"  ✓ Template: {template_path}")

    require_source(data, theme, skill_id=SKILL_ID)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Roleup: silently remove brown guide rectangles.
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # 1. Top placeholder (stella: main_message / roleup: chart_title)
    top_text = resolve_top_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:50]}")

    # 2. Subtitle placeholder (stella: chart_title / roleup: main_message)
    sub_text = resolve_subtitle_text(data, theme) or "対象会社概要：株主構成"
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:50]}")

    # 3. Remove template Table 1 (stella only; roleup template doesn't have it).
    template_table = find_shape(slide, SHAPE_TABLE, warn=False)
    if template_table is not None:
        slide.shapes._spTree.remove(template_table._element)
        print(f"  ✓ Template Table 1 removed")

    # 4. Format rows
    shareholders = data.get("shareholders", {})
    directors = data.get("directors", {})
    sh_rows = format_shareholder_rows(shareholders)
    dir_rows = format_director_rows(directors)

    has_total = shareholders.get("total") is not None
    n_data_rows = len(sh_rows) + len(dir_rows)
    n_headers = 2

    # 5. Layout calculation
    SLIDE_HEIGHT = prs.slide_height  # brand-aware: A4 横 for roleup, 16:9 for stella
    SECTION_TITLE_H = Inches(0.28)
    SECTION_GAP = Inches(0.12)

    total_available = SLIDE_HEIGHT - CONTENT_START_Y - BOTTOM_MARGIN
    overhead = (SECTION_TITLE_H * 2) + SECTION_GAP
    table_budget = total_available - overhead

    # 6. Pick font size
    chosen_font = None
    chosen_row_h = None

    for font_sz in FONT_CANDIDATES:
        row_h = required_row_height(font_sz)
        if row_h > MAX_ROW_HEIGHT:
            row_h = MAX_ROW_HEIGHT

        total_table_h = (HEADER_HEIGHT * n_headers) + (row_h * n_data_rows)

        if total_table_h <= table_budget:
            chosen_font = font_sz
            chosen_row_h = row_h
            print(f"  ✓ フォント {font_sz/100:.0f}pt: "
                  f"行高さ={row_h/914400:.2f}in, "
                  f"テーブル合計={total_table_h/914400:.2f}in ≤ "
                  f"予算{table_budget/914400:.2f}in → OK")
            break
        else:
            print(f"  × フォント {font_sz/100:.0f}pt: "
                  f"行高さ={row_h/914400:.2f}in, "
                  f"テーブル合計={total_table_h/914400:.2f}in > "
                  f"予算{table_budget/914400:.2f}in → 次へ")

    if chosen_font is None:
        chosen_font = FONT_CANDIDATES[-1]
        available_per_row = (table_budget - HEADER_HEIGHT * n_headers) / n_data_rows
        available_per_row = int(available_per_row)

        if available_per_row < MIN_ROW_HEIGHT:
            print(f"\n  ⚠ エラー: データ行数が多すぎます ({n_data_rows}行)")
            print(f"    最小フォント{FONT_CANDIDATES[-1]/100:.0f}ptでも "
                  f"行高さが{available_per_row/914400:.2f}inとなり、")
            print(f"    下限{MIN_ROW_HEIGHT/914400:.2f}inを下回ります。")
            print(f"    対応: 株主数または役員数を削減 / 2枚に分割")
            sys.exit(1)

        chosen_row_h = available_per_row
        print(f"  △ 最小フォント{chosen_font/100:.0f}ptで逆算: "
              f"行高さ={chosen_row_h/914400:.2f}in")

    font_size = chosen_font
    row_height = chosen_row_h

    # 7. Place sections
    current_y = CONTENT_START_Y

    # ■株主構成
    sh_section_title = shareholders.get("section_title", "株主構成")
    add_section_title(slide, f"■{sh_section_title}", TABLE_LEFT, current_y, TABLE_WIDTH)
    current_y += SECTION_TITLE_H

    sh_columns = shareholders.get("columns",
        ["#", "株主", "役職", "関係", "持株数", "議決権比率(%)", "備考"])
    # column widths scaled to TABLE_WIDTH (stella 12.52in → roleup 10.873in)
    width_in = TABLE_WIDTH / 914400
    sh_ratio = [0.50, 1.80, 1.30, 1.00, 1.20, 1.50, 5.22]
    sh_total = sum(sh_ratio)
    sh_col_widths = [w / sh_total * width_in for w in sh_ratio]

    sh_shape, sh_table_h = build_table(
        slide, sh_columns, sh_rows, sh_col_widths,
        TABLE_LEFT, current_y, TABLE_WIDTH,
        font_size=font_size, has_total=has_total,
        row_height=row_height,
    )
    current_y += sh_table_h + SECTION_GAP

    # ■役員構成
    dir_section_title = directors.get("section_title", "役員構成")
    add_section_title(slide, f"■{dir_section_title}", TABLE_LEFT, current_y, TABLE_WIDTH)
    current_y += SECTION_TITLE_H

    dir_columns = directors.get("columns",
        ["#", "氏名", "役職", "関係", "役員報酬", "備考"])
    dir_ratio = [0.50, 1.80, 1.30, 1.00, 1.50, 6.42]
    dir_total = sum(dir_ratio)
    dir_col_widths = [w / dir_total * width_in for w in dir_ratio]

    dir_shape, dir_table_h = build_table(
        slide, dir_columns, dir_rows, dir_col_widths,
        TABLE_LEFT, current_y, TABLE_WIDTH,
        font_size=font_size, has_total=False,
        row_height=row_height,
    )
    current_y += dir_table_h + Inches(0.05)

    # 8. Source
    source_text = data.get("source", "")
    if source_text:
        body = source_text if source_text.startswith("出典") else f"出典：{source_text}"
        source_shape = find_shape(slide, SHAPE_SOURCE, warn=False)
        if source_shape is not None:
            write_source_placeholder(source_shape, body)
            print(f"  ✓ 出典 ({SHAPE_SOURCE} placeholder): {body[:50]}")
        else:
            add_dynamic_source_textbox(slide, body)

    # 9. Save
    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
