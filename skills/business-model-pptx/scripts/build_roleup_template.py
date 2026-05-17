"""One-shot generator: business-model-pptx の roleup curated template を
stella template から派生して作る。

実行は本ファイルを直接 python3 で起動するだけ。生成済テンプレは
`assets/roleup/business-model-template.pptx` に上書き保存する。

設計:
  - スライドサイズ: A4 landscape (10691813 × 7559675 EMU = 11.69 × 8.27 in)
  - Title 1, Text Placeholder 2: roleup pilot (business-overview) の座標に揃える
  - Rectangle 4 (diagram area): 左 70% 幅 (~7.5 in)
  - TextBox 9 (implications): 右側 25% 幅 (~3.0 in)
  - Source 3: 下端のテキストボックス
  - フォント: Yu Gothic UI (Title / Subtitle / Source / Implications)
"""
import copy
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(SKILL_DIR, "assets", "stellar_aiz", "business-model-template.pptx")
DST = os.path.join(SKILL_DIR, "assets", "roleup", "business-model-template.pptx")

# A4 landscape EMU
A4_W = 10691813
A4_H = 7559675

FONT_EA = "Yu Gothic UI"
COLOR_TEXT = RGBColor(0x24, 0x1A, 0x17)
COLOR_SUBTITLE = RGBColor(0x89, 0x71, 0x41)
COLOR_SOURCE = RGBColor(0x3E, 0x3A, 0x39)


def _find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def _set_run_font(run, font_name, color, size_pt=None, bold=False):
    run.font.name = font_name
    run.font.color.rgb = color
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    run.font.bold = bold
    # ensure latin/ea both
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn("a:rPr"))
    for tag in ("a:latin", "a:ea"):
        el = rPr.find(qn(tag))
        if el is not None:
            rPr.remove(el)
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": font_name})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": font_name})


def main():
    prs = Presentation(SRC)
    prs.slide_width = A4_W
    prs.slide_height = A4_H

    slide = prs.slides[0]

    # Remove think-cell OLE if still present (silent)
    for s in list(slide.shapes):
        if "think-cell" in (s.name or ""):
            sp = s._element
            sp.getparent().remove(sp)

    # Title 1
    title = _find(slide, "Title 1")
    if title is not None:
        title.left = Inches(0.41)
        title.top = Inches(0.41)
        title.width = Inches(10.87)
        title.height = Inches(0.50)
        if title.has_text_frame:
            for para in title.text_frame.paragraphs:
                for run in para.runs:
                    _set_run_font(run, FONT_EA, COLOR_TEXT, size_pt=22, bold=True)

    # Text Placeholder 2 (Chart Title / Subtitle)
    sub = _find(slide, "Text Placeholder 2")
    if sub is not None:
        sub.left = Inches(0.41)
        sub.top = Inches(1.00)
        sub.width = Inches(10.87)
        sub.height = Inches(0.36)
        if sub.has_text_frame:
            for para in sub.text_frame.paragraphs:
                for run in para.runs:
                    _set_run_font(run, FONT_EA, COLOR_SUBTITLE, size_pt=12, bold=False)

    # Rectangle 4 (diagram area) — 左 70%
    rect = _find(slide, "Rectangle 4")
    if rect is not None:
        rect.left = Inches(0.41)
        rect.top = Inches(1.55)
        rect.width = Inches(7.40)
        rect.height = Inches(5.40)

    # TextBox 9 (implications) — 右側
    tb9 = _find(slide, "TextBox 9")
    if tb9 is not None:
        tb9.left = Inches(8.00)
        tb9.top = Inches(1.55)
        tb9.width = Inches(3.28)
        tb9.height = Inches(5.40)
        if tb9.has_text_frame:
            for pi, para in enumerate(tb9.text_frame.paragraphs):
                # Para[0]: 「意味合い」見出し
                # Para[1..3]: implication items
                for run in para.runs:
                    if pi == 0:
                        _set_run_font(run, FONT_EA, COLOR_TEXT, size_pt=12, bold=True)
                    else:
                        _set_run_font(run, FONT_EA, COLOR_TEXT, size_pt=10, bold=False)

    # Source 3 — 下端に新規追加 (TextBox)
    if _find(slide, "Source 3") is None:
        src_tb = slide.shapes.add_textbox(
            Inches(0.41), Inches(7.30), Inches(10.87), Inches(0.40)
        )
        src_tb.name = "Source 3"
        tf = src_tb.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "出典: "
        _set_run_font(run, FONT_EA, COLOR_SOURCE, size_pt=6, bold=False)

    prs.save(DST)
    print(f"✅ Generated: {DST}")
    print(f"   slide_size: {prs.slide_width} × {prs.slide_height} EMU")
    for s in slide.shapes:
        print(f"   - {s.name}")


if __name__ == "__main__":
    main()
