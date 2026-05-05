"""
fill_pest_analysis.py — PEST分析スライドをPPTXネイティブオブジェクトで生成

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 2×2 PESTマトリクス:
      [P: 政治要因 ] [E: 経済要因 ]   ← 上段
      [S: 社会要因 ] [T: 技術要因 ]   ← 下段
  - 下部: 出典

各象限:
  - 上部にカラーヘッダーバー（象限ラベル + 英語名）
  - 下部にブレット項目リスト（各項目に影響度インジケーター ↑/↓/→ オプション）

Usage:
  python fill_pest_analysis.py \
    --data /home/claude/pest_data.json \
    --template <path>/pest-analysis-template.pptx \
    --output /mnt/user-data/outputs/PEST_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "pest-analysis-pptx"

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Layout Constants (defaults; reassigned by _apply_theme) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE = "Source 3"

GRID_LEFT = Inches(0.41)
GRID_TOP = Inches(1.55)
GRID_WIDTH = Inches(12.51)
GRID_HEIGHT = Inches(5.30)

GAP = Inches(0.15)
HEADER_H = Inches(0.55)

CELL_W = (GRID_WIDTH - GAP) / 2
CELL_H = (GRID_HEIGHT - GAP) / 2

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(6.93)
SOURCE_W = Inches(12.50)
SOURCE_H = Inches(0.30)

# ── Colors ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ─── 共通配色（正本: skills/_common/styles/chart_palette.md） ───
# 編集時は _common/styles/chart_palette.md と他 4 スキルの fill_*.py も同期更新
# CHART_PALETTE には TARGET_COLOR(赤) と OTHER_COLOR(灰) を含めない（palette 外で固定）
CHART_PALETTE = [
    "#4E79A7", "#F28E2B", "#59A14F", "#76B7B2",
    "#EDC948", "#B07AA1", "#FF9DA7", "#9C755F",
]
OTHER_COLOR = "#BAB0AC"
TARGET_COLOR = "#E15759"
LABEL_BAR_COLOR = "#4E79A7"
LABEL_BG_COLOR = "#E8EEF5"


def _palette_color(index: int, total: int) -> str:
    if total <= 1:
        return CHART_PALETTE[0]
    return CHART_PALETTE[index % len(CHART_PALETTE)]

# PEST 各象限の色
# 新仕様（2026-04-29）: PEST 4象限は全て単色青に統一（CATEGORY_COLORS と同じ思想）
# 旧色（紫/青/緑/橙）は LABEL_BAR_RGB に統一
LABEL_BAR_RGB = RGBColor(0x4E, 0x79, 0xA7)
LABEL_BG_RGB = RGBColor(0xE8, 0xEE, 0xF5)
COLOR_P = LABEL_BAR_RGB
COLOR_E = LABEL_BAR_RGB
COLOR_S = LABEL_BAR_RGB
COLOR_T = LABEL_BAR_RGB

# 薄い背景色
# 新仕様: 4象限の塗り背景も単色（薄い青 LABEL_BG_RGB）に統一
COLOR_P_LIGHT = LABEL_BG_RGB
COLOR_E_LIGHT = LABEL_BG_RGB
COLOR_S_LIGHT = LABEL_BG_RGB
COLOR_T_LIGHT = LABEL_BG_RGB

# 影響度インジケーター色
COLOR_IMPACT_POSITIVE = RGBColor(0x1B, 0x7A, 0x3B)    # 濃緑
COLOR_IMPACT_NEGATIVE = RGBColor(0xC0, 0x3A, 0x3A)    # 濃赤
COLOR_IMPACT_NEUTRAL = RGBColor(0x66, 0x66, 0x66)     # グレー

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_HEADER = Pt(16)
FONT_SIZE_HEADER_LETTER = Pt(20)   # ヘッダーバーの大きな P/E/S/T 字 (stella: 16×1.2)
FONT_SIZE_HEADER_EN = Pt(11)
FONT_SIZE_ITEM = Pt(12)
FONT_SIZE_SOURCE = Pt(10)
FONT_SIZE_LEGEND = Pt(10)

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global GRID_LEFT, GRID_TOP, GRID_WIDTH, GRID_HEIGHT, GAP, HEADER_H, CELL_W, CELL_H
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE, LABEL_BAR_RGB, LABEL_BG_RGB
    global COLOR_P, COLOR_E, COLOR_S, COLOR_T
    global COLOR_P_LIGHT, COLOR_E_LIGHT, COLOR_S_LIGHT, COLOR_T_LIGHT
    global FONT_NAME_JP
    global FONT_SIZE_HEADER, FONT_SIZE_HEADER_LETTER, FONT_SIZE_HEADER_EN
    global FONT_SIZE_ITEM, FONT_SIZE_SOURCE, FONT_SIZE_LEGEND

    _THEME = theme

    GRID_LEFT = theme.layout("grid_left_in")
    GRID_TOP = theme.layout("grid_top_in")
    GRID_WIDTH = theme.layout("grid_width_in")
    GRID_HEIGHT = theme.layout("grid_height_in")
    GAP = theme.layout("gap_in")
    HEADER_H = theme.layout("header_h_in")
    CELL_W = (GRID_WIDTH - GAP) / 2
    CELL_H = (GRID_HEIGHT - GAP) / 2
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    if theme.id == "roleup":
        # Brown PEST palette (single-tone like stella's blue unification).
        LABEL_BAR_RGB = theme.color("label_bar")
        LABEL_BG_RGB = _hex_to_rgbcolor(theme.hex("label_bg"))
        FONT_NAME_JP = theme.font_ea
        # roleup C4 [22,14,12,10,6]
        FONT_SIZE_HEADER = theme.pt("font_size_subtitle_pt")        # 12pt
        FONT_SIZE_HEADER_LETTER = theme.pt("font_size_key_message_pt")  # 14pt
        FONT_SIZE_HEADER_EN = theme.font_size_body_pt(skill_id=SKILL_ID)  # 10pt
        FONT_SIZE_ITEM = theme.pt("font_size_subtitle_pt")          # 12pt
        FONT_SIZE_SOURCE = theme.pt("font_size_source_pt")          # 6pt
        FONT_SIZE_LEGEND = theme.font_size_body_pt(skill_id=SKILL_ID)  # 10pt
    # else stella: keep V1 hardcoded values for regression-zero.

    COLOR_P = LABEL_BAR_RGB
    COLOR_E = LABEL_BAR_RGB
    COLOR_S = LABEL_BAR_RGB
    COLOR_T = LABEL_BAR_RGB
    COLOR_P_LIGHT = LABEL_BG_RGB
    COLOR_E_LIGHT = LABEL_BG_RGB
    COLOR_S_LIGHT = LABEL_BG_RGB
    COLOR_T_LIGHT = LABEL_BG_RGB


def _hex_to_rgbcolor(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _silent_remove_shape(slide, shape_name: str) -> None:
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


# ──────────────────────────────────────────────
# Utility
# ──────────────────────────────────────────────
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = FONT_NAME_JP
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


def _parse_impact(item):
    """
    アイテムを {text, impact} の形に正規化する。
    - str の場合: {"text": str, "impact": None}
    - dict の場合: そのまま
    """
    if isinstance(item, str):
        return {"text": item, "impact": None}
    elif isinstance(item, dict):
        return {
            "text": item.get("text", ""),
            "impact": item.get("impact"),
        }
    return {"text": str(item), "impact": None}


def _impact_char_and_color(impact):
    """impact文字列 → (記号, RGB色) を返す"""
    if impact is None:
        return (None, None)
    impact_lower = str(impact).lower()
    if impact_lower in ("positive", "+", "up", "追い風", "pos"):
        return ("▲", COLOR_IMPACT_POSITIVE)
    elif impact_lower in ("negative", "-", "down", "逆風", "neg"):
        return ("▼", COLOR_IMPACT_NEGATIVE)
    elif impact_lower in ("neutral", "=", "flat", "中立", "neu"):
        return ("▬", COLOR_IMPACT_NEUTRAL)
    return (None, None)


# ──────────────────────────────────────────────
# PEST Quadrant Builder
# ──────────────────────────────────────────────
def build_quadrant(slide, label_jp, label_en, label_letter, items,
                    header_color, body_color,
                    left, top, width, height):
    """
    1つの象限を描画する。
      - 上部: カラーヘッダーバー（文字ラベル + 日本語ラベル + 英語名）
      - 下部: 薄い背景 + ブレット項目リスト（影響度インジケーター付き）
    """
    # 全体の外枠
    outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    outer.fill.solid()
    outer.fill.fore_color.rgb = body_color
    outer.line.color.rgb = header_color
    outer.line.width = Pt(0.75)
    outer.shadow.inherit = False
    outer.text_frame.text = ""

    # ヘッダーバー
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, HEADER_H
    )
    header.fill.solid()
    header.fill.fore_color.rgb = header_color
    header.line.fill.background()
    header.shadow.inherit = False

    tf = header.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 既存段落クリア
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    p = tf._txBody.makeelement(qn("a:p"), {})
    tf._txBody.append(p)
    pPr = etree.SubElement(p, qn("a:pPr"))
    pPr.set("algn", "l")

    # Run 1: 文字ラベル（P/E/S/T） 大きく
    r0 = etree.SubElement(p, qn("a:r"))
    rPr0 = etree.SubElement(r0, qn("a:rPr"), attrib={
        "lang": "en-US",
        "sz": str(int(FONT_SIZE_HEADER_LETTER.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr0, qn("a:latin"), attrib={"typeface": FONT_NAME_JP if (_THEME and _THEME.id == "roleup") else "Arial"})
    sf0 = etree.SubElement(rPr0, qn("a:solidFill"))
    s0 = etree.SubElement(sf0, qn("a:srgbClr"))
    s0.set("val", "FFFFFF")
    t0 = etree.SubElement(r0, qn("a:t"))
    t0.text = f"{label_letter}  "

    # Run 2: 日本語ラベル
    r1 = etree.SubElement(p, qn("a:r"))
    rPr1 = etree.SubElement(r1, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(FONT_SIZE_HEADER.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr1, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr1, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    sf1 = etree.SubElement(rPr1, qn("a:solidFill"))
    s1 = etree.SubElement(sf1, qn("a:srgbClr"))
    s1.set("val", "FFFFFF")
    t1 = etree.SubElement(r1, qn("a:t"))
    t1.text = f"{label_jp}  "

    # Run 3: 英語ラベル
    r2 = etree.SubElement(p, qn("a:r"))
    rPr2 = etree.SubElement(r2, qn("a:rPr"), attrib={
        "lang": "en-US",
        "sz": str(int(FONT_SIZE_HEADER_EN.pt * 100)),
        "b": "0",
    })
    etree.SubElement(rPr2, qn("a:latin"), attrib={"typeface": FONT_NAME_JP if (_THEME and _THEME.id == "roleup") else "Arial"})
    sf2 = etree.SubElement(rPr2, qn("a:solidFill"))
    s2 = etree.SubElement(sf2, qn("a:srgbClr"))
    s2.set("val", "FFFFFF")
    t2 = etree.SubElement(r2, qn("a:t"))
    t2.text = f"({label_en})"

    # ボディ部分
    body_top = top + HEADER_H
    body_h = height - HEADER_H

    body_box = slide.shapes.add_textbox(
        left + Inches(0.15), body_top + Inches(0.12),
        width - Inches(0.30), body_h - Inches(0.20),
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.margin_left = 0; body_tf.margin_right = 0
    body_tf.margin_top = 0; body_tf.margin_bottom = 0
    body_tf.vertical_anchor = MSO_ANCHOR.TOP

    # 既存段落クリア
    for p in list(body_tf.paragraphs):
        p._p.getparent().remove(p._p)

    header_color_hex = "{:02X}{:02X}{:02X}".format(
        header_color[0], header_color[1], header_color[2]
    )

    for i, raw_item in enumerate(items):
        item = _parse_impact(raw_item)
        text = item["text"]
        impact = item["impact"]
        impact_char, impact_rgb = _impact_char_and_color(impact)

        p_elem = etree.SubElement(body_tf._txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"), attrib={
            "marL": "180000",
            "indent": "-180000",
        })
        if i > 0:
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "400"})

        # Bullet character
        buChar = etree.SubElement(pPr, qn("a:buChar"), attrib={"char": "•"})
        buFont = etree.SubElement(pPr, qn("a:buFont"), attrib={"typeface": "Arial"})
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        buClrSolid = etree.SubElement(buClr, qn("a:srgbClr"))
        buClrSolid.set("val", header_color_hex)

        # 影響度インジケーター（あれば）
        if impact_char and impact_rgb:
            # Run: インジケーター記号
            r_imp = etree.SubElement(p_elem, qn("a:r"))
            rPr_imp = etree.SubElement(r_imp, qn("a:rPr"), attrib={
                "lang": "en-US",
                "sz": str(int(FONT_SIZE_ITEM.pt * 100)),
                "b": "1",
            })
            etree.SubElement(rPr_imp, qn("a:latin"), attrib={"typeface": FONT_NAME_JP if (_THEME and _THEME.id == "roleup") else "Arial"})
            sf_imp = etree.SubElement(rPr_imp, qn("a:solidFill"))
            s_imp = etree.SubElement(sf_imp, qn("a:srgbClr"))
            s_imp.set("val", "{:02X}{:02X}{:02X}".format(
                impact_rgb[0], impact_rgb[1], impact_rgb[2]
            ))
            t_imp = etree.SubElement(r_imp, qn("a:t"))
            t_imp.text = f"{impact_char} "

        # Run: テキスト本文
        r = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(FONT_SIZE_ITEM.pt * 100)),
        })
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        s = etree.SubElement(sf, qn("a:srgbClr"))
        s.set("val", "333333")
        t = etree.SubElement(r, qn("a:t"))
        t.text = text

    print(f"  ✓ 象限 [{label_letter}: {label_jp}]: {len(items)}項目")


# ──────────────────────────────────────────────
# Impact Legend (bottom left, optional)
# ──────────────────────────────────────────────
def add_impact_legend(slide):
    """画面下部に影響度の凡例を追加（▲追い風 / ▬中立 / ▼逆風）"""
    legend_y = Inches(6.93)
    # "影響度:" ラベル + 3つの記号を右詰めで配置
    # 左下角は出典に使うので、凡例はもう少し上、または右寄せにする
    # 今回はシンプルに出典と同じY座標の右側に配置
    legend_items = [
        ("▲ 追い風", COLOR_IMPACT_POSITIVE),
        ("▬ 中立", COLOR_IMPACT_NEUTRAL),
        ("▼ 逆風", COLOR_IMPACT_NEGATIVE),
    ]
    legend_x = Inches(9.50)
    legend_item_w = Inches(1.10)

    for i, (text, color) in enumerate(legend_items):
        tb = slide.shapes.add_textbox(
            legend_x + legend_item_w * i,
            legend_y,
            legend_item_w, Inches(0.25),
        )
        tf = tb.text_frame
        tf.margin_left = 0; tf.margin_right = 0
        tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        # 記号を色付きで表示
        symbol = text.split()[0]  # ▲/▬/▼
        label = text.split()[1]   # 追い風/中立/逆風

        r1 = p.add_run()
        r1.text = f"{symbol} "
        r1.font.size = FONT_SIZE_LEGEND
        r1.font.bold = True
        r1.font.name = "Arial"
        r1.font.color.rgb = color

        r2 = p.add_run()
        r2.text = label
        r2.font.size = FONT_SIZE_LEGEND
        r2.font.bold = False
        r2.font.name = FONT_NAME_JP
        r2.font.color.rgb = COLOR_SOURCE


# ──────────────────────────────────────────────
# Validation
# ──────────────────────────────────────────────
MAIN_MESSAGE_MAX = 65
PEST_ITEM_MAX = 50


def _validate_input(data):
    """main_message ≤65字、各 quadrant の item.text ≤50字。"""
    main_message = data.get("main_message", "")
    if len(main_message) > MAIN_MESSAGE_MAX:
        raise ValueError(
            f"main_message は {MAIN_MESSAGE_MAX} 字以内（受領: {len(main_message)}）: {main_message[:80]}..."
        )
    pest = data.get("pest", {})
    for axis in ("political", "economic", "social", "technological"):
        items = pest.get(axis, {}).get("items", [])
        for j, it in enumerate(items):
            text = it.get("text", "")
            if len(text) > PEST_ITEM_MAX:
                raise ValueError(
                    f"pest.{axis}.items[{j}].text は {PEST_ITEM_MAX} 字以内"
                    f"（受領: {len(text)}）: {text}"
                )


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand.",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "pest-analysis")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    require_source(data, theme, skill_id=SKILL_ID)
    _validate_input(data)
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top / subtitle placeholder (brand 別)
    top_text = resolve_top_text(data, theme)
    sub_text = resolve_subtitle_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE),
                     sub_text or data.get("chart_title", "PEST分析"))
    print(f"  ✓ Main Message & Chart Title set")

    # Roleup: silently remove brown guide rectangles carried by cp-derived template.
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    pest = data.get("pest", {})

    # 4象限の設定
    quadrants = [
        {
            "data": pest.get("political", {}),
            "label_jp": "政治要因",
            "label_en": "Political",
            "label_letter": "P",
            "header_color": COLOR_P,
            "body_color": COLOR_P_LIGHT,
            "left": GRID_LEFT,
            "top": GRID_TOP,
        },
        {
            "data": pest.get("economic", {}),
            "label_jp": "経済要因",
            "label_en": "Economic",
            "label_letter": "E",
            "header_color": COLOR_E,
            "body_color": COLOR_E_LIGHT,
            "left": GRID_LEFT + CELL_W + GAP,
            "top": GRID_TOP,
        },
        {
            "data": pest.get("social", {}),
            "label_jp": "社会要因",
            "label_en": "Social",
            "label_letter": "S",
            "header_color": COLOR_S,
            "body_color": COLOR_S_LIGHT,
            "left": GRID_LEFT,
            "top": GRID_TOP + CELL_H + GAP,
        },
        {
            "data": pest.get("technological", {}),
            "label_jp": "技術要因",
            "label_en": "Technological",
            "label_letter": "T",
            "header_color": COLOR_T,
            "body_color": COLOR_T_LIGHT,
            "left": GRID_LEFT + CELL_W + GAP,
            "top": GRID_TOP + CELL_H + GAP,
        },
    ]

    # 影響度インジケーターが1つでも使われているか
    has_impact = False
    for q in quadrants:
        items = q["data"].get("items", [])
        for itm in items:
            if isinstance(itm, dict) and itm.get("impact"):
                has_impact = True
                break
        if has_impact:
            break

    for q in quadrants:
        label_jp = q["data"].get("label_jp", q["label_jp"])
        label_en = q["data"].get("label_en", q["label_en"])
        items = q["data"].get("items", [])
        build_quadrant(
            slide, label_jp, label_en, q["label_letter"], items,
            q["header_color"], q["body_color"],
            q["left"], q["top"], CELL_W, CELL_H,
        )

    # 凡例（影響度インジケーターが使われている場合のみ）
    # roleup は A4 横でグリッドが下方まで広がるため、凡例の配置スペースが無く、
    # ▲/▬/▼ 記号の色 (緑/灰/赤) で直感的に意味が伝わる前提で凡例を省略。
    if has_impact and theme.id == "stellar_aiz":
        add_impact_legend(slide)
        print(f"  ✓ 影響度凡例を追加")

    # 出典: roleup は Source 3 placeholder にセット、stella は dynamic textbox を追加
    source = data.get("source", "")
    if source:
        if theme.id == "stellar_aiz":
            # 影響度凡例 (Inches(9.50)〜) と重ならないよう幅を制限
            add_text_box(
                slide, source,
                SOURCE_X, SOURCE_Y, Inches(8.90), SOURCE_H,
                FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                align=PP_ALIGN.LEFT,
            )
        else:
            source_shape = find_shape(slide, SHAPE_SOURCE)
            if source_shape is not None:
                set_textbox_text(source_shape, source)
            else:
                add_text_box(
                    slide, source,
                    SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                    FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                    align=PP_ALIGN.LEFT,
                )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
