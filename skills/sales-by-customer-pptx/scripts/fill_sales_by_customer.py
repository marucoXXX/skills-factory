"""
fill_sales_by_customer.py — 主要販売先売上高テーブルをPPTXネイティブテーブルで生成

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。

レイアウト:
  - 上部: メインメッセージ (Title 1) + チャートタイトル (Text Placeholder 2)
  - 中央: 期数 N 個の独立テーブルを横並び (#, 企業名, 売上高, 割合)
  - 下部: 出典 (stella=Source textbox / roleup=Source 3 placeholder)

継続顧客 (全期に出現する企業) は強調色で色付け。

Usage:
  python fill_sales_by_customer.py --brand stellar_aiz \\
    --data {{WORK_DIR}}/sales_by_customer_data.json \\
    --output {{OUTPUT_DIR}}/SalesByCustomer_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "sales-by-customer-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names (template-structure invariants) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_CONTENT_AREA = "Content Area"

# Defaults (stella). Reassigned in _apply_theme(theme).
SHAPE_SOURCE = "Source"

CONTENT_LEFT = Inches(0.41)
CONTENT_TOP = Inches(1.50)
CONTENT_WIDTH = Inches(12.52)
CONTENT_BOTTOM = Inches(6.90)
TABLE_GAP = Inches(0.15)
NOTE_HEIGHT = Inches(0.25)
PERIOD_HEADER_HEIGHT = Inches(0.30)
BOTTOM_MARGIN = Inches(0.05)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(8.00)
SOURCE_H = Inches(0.30)

COL_RATIOS = [0.08, 0.42, 0.30, 0.20]  # #, 企業名, 売上高, 割合

# Stella defaults; reassigned in _apply_theme(theme).
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_HIGHLIGHT = RGBColor(0xE6, 0x7E, 0x00)  # 継続顧客 (継続的な取引)
COLOR_HEADER_BG = RGBColor(0xF0, 0xF0, 0xF0)
COLOR_EVEN_ROW = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_OTHER_ROW = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BORDER_HEX = "E0E0E0"
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)

FONT_NAME = "Meiryo UI"
FONT_NAME_LATIN = "Arial"
SOURCE_FONT_PT = 10

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global SHAPE_SOURCE
    global CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_BOTTOM
    global TABLE_GAP, NOTE_HEIGHT, PERIOD_HEADER_HEIGHT, BOTTOM_MARGIN
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_HIGHLIGHT, COLOR_HEADER_BG, COLOR_EVEN_ROW, COLOR_OTHER_ROW
    global COLOR_BORDER_HEX, COLOR_SOURCE
    global FONT_NAME, FONT_NAME_LATIN, SOURCE_FONT_PT

    _THEME = theme
    FONT_NAME = theme.font_ea
    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    CONTENT_LEFT = theme.layout("content_left_in")
    CONTENT_TOP = theme.layout("content_top_in")
    CONTENT_WIDTH = theme.layout("content_width_in")
    CONTENT_BOTTOM = theme.layout("content_bottom_in")
    TABLE_GAP = theme.layout("table_gap_in")
    NOTE_HEIGHT = theme.layout("note_height_in")
    PERIOD_HEADER_HEIGHT = theme.layout("period_header_height_in")
    BOTTOM_MARGIN = theme.layout("bottom_margin_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        FONT_NAME_LATIN = "Arial"
        COLOR_HIGHLIGHT = RGBColor(0xE6, 0x7E, 0x00)
        COLOR_HEADER_BG = RGBColor(0xF0, 0xF0, 0xF0)
        COLOR_EVEN_ROW = RGBColor(0xFA, 0xFA, 0xFA)
        COLOR_OTHER_ROW = RGBColor(0xF5, 0xF5, 0xF5)
        COLOR_BORDER_HEX = "E0E0E0"
        SOURCE_FONT_PT = 10
    else:
        SHAPE_SOURCE = "Source 3"
        FONT_NAME_LATIN = theme.font_ea  # roleup: Yu Gothic UI for both latin & ea
        COLOR_HIGHLIGHT = theme.color("highlight_target")     # #C78624
        COLOR_HEADER_BG = theme.color("header_bg")            # #F5EFE5
        COLOR_EVEN_ROW = theme.color("label_bg")              # #F2E8DD (banded)
        COLOR_OTHER_ROW = theme.color("label_bg")             # #F2E8DD (その他行)
        COLOR_BORDER_HEX = theme.hex_no_hash("highlight_other")  # #CDCECE
        SOURCE_FONT_PT = theme.pt_value("font_size_source_pt")  # 6


def _silent_remove_shape(slide, name):
    for shape in list(slide.shapes):
        if shape.name == name:
            slide.shapes._spTree.remove(shape._element)
            return True
    return False


def find_shape(slide, name, warn=True):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    if warn:
        print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def remove_shape(slide, name):
    shape = find_shape(slide, name, warn=False)
    if shape is not None:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)
        print(f"  ✓ Shape '{name}' removed")


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def get_font_sizes(n_periods):
    """期数に応じたフォントサイズを返す。

    stella: 期数によって 16/14/13 等の dynamic スケール (regression-zero)。
    roleup: C4 許容セット {22, 14, 12, 10, 6} に固定。期数に依らず本文 10pt。
    """
    if _THEME is not None and _THEME.id != "stellar_aiz":
        # roleup: 全期数で同じ。C4 許容セット内のみ。
        return {
            "header": Pt(_THEME.pt_value("font_size_subtitle_pt")),       # 12pt
            "body": Pt(_THEME.pt_value("font_size_body_pt")),             # 10pt
            "period": Pt(_THEME.pt_value("font_size_key_message_pt")),    # 14pt
            "note": Pt(_THEME.pt_value("font_size_body_pt")),             # 10pt
        }
    # stella V1 互換
    if n_periods <= 3:
        return {"header": Pt(14), "body": Pt(13), "period": Pt(16), "note": Pt(11)}
    elif n_periods == 4:
        return {"header": Pt(12), "body": Pt(11), "period": Pt(14), "note": Pt(10)}
    else:  # 5
        return {"header": Pt(11), "body": Pt(10), "period": Pt(13), "note": Pt(9)}


def _num_to_kanji(n):
    return {2: "二", 3: "三", 4: "四", 5: "五"}.get(n, str(n))


def identify_continuous_customers(periods):
    """全期間に出現する企業名を特定"""
    if len(periods) < 2:
        return set()
    name_sets = []
    for period in periods:
        names = set()
        for row in period.get("customers", []):
            name = row.get("name", "")
            if name and name != "その他":
                names.add(name)
        name_sets.append(names)
    result = name_sets[0]
    for ns in name_sets[1:]:
        result = result & ns
    return result


def _color_hex_no_hash(rgb):
    if isinstance(rgb, RGBColor):
        return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    if isinstance(rgb, (tuple, list)):
        return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    return str(rgb).replace("#", "")


def add_textbox(slide, left, top, width, height, text, font_size,
                font_bold=False, font_color=None, alignment=PP_ALIGN.LEFT,
                font_underline=False):
    if font_color is None:
        font_color = COLOR_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = font_size
    run.font.bold = font_bold
    run.font.color.rgb = font_color
    if font_underline:
        run.font.underline = True
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", FONT_NAME_LATIN)
    return txBox


def set_cell_text(cell, text, font_size, font_bold=False, font_color=None,
                  alignment=PP_ALIGN.LEFT, bg_color=None):
    if font_color is None:
        font_color = COLOR_TEXT

    tc = cell._tc

    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        txBody = etree.SubElement(tc, qn("a:txBody"))

    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
    bodyPr.set("lIns", "45720")
    bodyPr.set("rIns", "45720")
    bodyPr.set("tIns", "18288")
    bodyPr.set("bIns", "18288")
    bodyPr.set("anchor", "ctr")

    if txBody.find(qn("a:lstStyle")) is None:
        etree.SubElement(txBody, qn("a:lstStyle"))

    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    p_elem = etree.SubElement(txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    align_map = {PP_ALIGN.LEFT: "l", PP_ALIGN.CENTER: "ctr", PP_ALIGN.RIGHT: "r"}
    pPr.set("algn", align_map.get(alignment, "l"))

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(font_size.pt * 100)),
        "b": "1" if font_bold else "0",
        "dirty": "0",
    })
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", _color_hex_no_hash(font_color))
    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", FONT_NAME_LATIN)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT_NAME)

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = str(text)

    old_tcPr = tc.find(qn("a:tcPr"))
    if old_tcPr is not None:
        tc.remove(old_tcPr)

    tcPr = etree.SubElement(tc, qn("a:tcPr"))
    tcPr.set("marL", "45720")
    tcPr.set("marR", "45720")
    tcPr.set("marT", "18288")
    tcPr.set("marB", "18288")

    for border_name in ["a:lnL", "a:lnR", "a:lnT"]:
        ln = etree.SubElement(tcPr, qn(border_name))
        ln.set("w", "0")
        etree.SubElement(ln, qn("a:noFill"))

    lnB = etree.SubElement(tcPr, qn("a:lnB"))
    lnB.set("w", "6350")
    sfB = etree.SubElement(lnB, qn("a:solidFill"))
    cB = etree.SubElement(sfB, qn("a:srgbClr"))
    cB.set("val", COLOR_BORDER_HEX)

    if bg_color is not None:
        fill = etree.SubElement(tcPr, qn("a:solidFill"))
        clr = etree.SubElement(fill, qn("a:srgbClr"))
        clr.set("val", _color_hex_no_hash(bg_color))
    else:
        etree.SubElement(tcPr, qn("a:noFill"))


def build_period_table(slide, data_period, continuous_names, font_sizes,
                       tbl_left, tbl_top, tbl_width, tbl_height):
    customers = data_period.get("customers", [])
    unit = data_period.get("unit", "千円")
    n_rows = len(customers) + 1
    n_cols = 4

    shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = shape.table

    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    for ci, ratio in enumerate(COL_RATIOS):
        table.columns[ci].width = int(tbl_width * ratio)

    header_h = int(tbl_height * 0.09)
    data_h = int((tbl_height - header_h) / len(customers))

    for i, tr in enumerate(tbl_elem.findall(qn('a:tr'))):
        if i == 0:
            tr.set('h', str(header_h))
        else:
            tr.set('h', str(data_h))

    headers = ["#", "企業名", f"売上高({unit})", "割合"]
    aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]

    for ci, (hdr, align) in enumerate(zip(headers, aligns)):
        set_cell_text(table.cell(0, ci), hdr,
                      font_size=font_sizes["header"],
                      font_bold=True,
                      font_color=COLOR_TEXT,
                      alignment=align,
                      bg_color=COLOR_HEADER_BG)

    for ri, cust in enumerate(customers):
        name = cust.get("name", "")
        revenue = cust.get("revenue", 0)
        share = cust.get("share", 0)
        rank = cust.get("rank", ri + 1)

        is_continuous = name in continuous_names and name != "その他"
        font_color = COLOR_HIGHLIGHT if is_continuous else COLOR_TEXT
        font_bold = is_continuous

        rank_str = str(rank) if name != "その他" else ""
        rev_str = f"{int(revenue):,}" if isinstance(revenue, (int, float)) else str(revenue)
        share_str = f"{share:.1f}%" if isinstance(share, (int, float)) else str(share)

        if name == "その他":
            bg = COLOR_OTHER_ROW
        elif ri % 2 == 1:
            bg = COLOR_EVEN_ROW
        else:
            bg = COLOR_WHITE

        row_idx = ri + 1
        set_cell_text(table.cell(row_idx, 0), rank_str,
                      font_sizes["body"], font_bold, font_color, PP_ALIGN.CENTER, bg)
        set_cell_text(table.cell(row_idx, 1), name,
                      font_sizes["body"], font_bold, font_color, PP_ALIGN.LEFT, bg)
        set_cell_text(table.cell(row_idx, 2), rev_str,
                      font_sizes["body"], font_bold, font_color, PP_ALIGN.RIGHT, bg)
        set_cell_text(table.cell(row_idx, 3), share_str,
                      font_sizes["body"], font_bold, font_color, PP_ALIGN.RIGHT, bg)

    return shape


def add_source_label(slide, text):
    """出典: roleup なら Source 3 placeholder、stella なら textbox。"""
    src_shape = find_shape(slide, SHAPE_SOURCE, warn=False)
    if src_shape is not None and src_shape.has_text_frame:
        tf = src_shape.text_frame
        tf.word_wrap = True
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)
        p = tf.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(SOURCE_FONT_PT)
        run.font.color.rgb = COLOR_SOURCE
        run.font.name = FONT_NAME
        print(f"  ✓ 出典 ({SHAPE_SOURCE} placeholder): {text[:50]}")
        return

    txBox = slide.shapes.add_textbox(SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(SOURCE_FONT_PT)
    run.font.color.rgb = COLOR_SOURCE
    run.font.name = FONT_NAME
    print(f"  ✓ 出典 (textbox): {text[:50]}")


def main():
    parser = argparse.ArgumentParser(description="主要販売先売上高スライド生成（ネイティブテーブル方式）")
    parser.add_argument("--data", required=True, help="JSONデータファイルパス")
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand.",
    )
    parser.add_argument("--output", required=True, help="出力PPTXファイルパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "sales-by-customer")

    print(f"=== 主要販売先売上高スライド生成 (brand={theme.id}) ===")
    print(f"  Template: {template_path}")

    require_source(data, theme, skill_id=SKILL_ID)

    # chart_title のデフォルトを data に埋めて、brand 別の top/subtitle 解決で
    # どちらの placeholder field に書かれても落ちない状態にする。
    if not data.get("chart_title"):
        data["chart_title"] = "主要販売先からの完成工事売上高と割合"

    prs = Presentation(template_path)
    slide = prs.slides[0]

    periods = data.get("periods", [])
    n_periods = len(periods)
    font_sizes = get_font_sizes(n_periods)
    continuous_names = identify_continuous_customers(periods)

    print(f"  期数: {n_periods}")
    print(f"  継続顧客: {sorted(continuous_names)}")

    # Top placeholder (brand-aware)
    top_text = resolve_top_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:50]}")

    # Subtitle placeholder
    sub_text = resolve_subtitle_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:50]}")

    # Content Area 削除 (silent for roleup which lacks it)
    _silent_remove_shape(slide, SHAPE_CONTENT_AREA)

    # Roleup: silently remove brown guide rectangles
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # 注記テキストボックス (継続顧客の凡例)
    subtitle_note = data.get("subtitle_note", "")
    continuous_label = data.get("continuous_label", "継続的な顧客")
    continuous_threshold_label = data.get("continuous_threshold_label", "")

    if not subtitle_note and continuous_names:
        if continuous_threshold_label:
            color_word = "継続色" if (_THEME is not None and _THEME.id != "stellar_aiz") else "オレンジ字"
            subtitle_note = f"{color_word}：{continuous_label}（{continuous_threshold_label}）"
        else:
            color_word = "継続色" if (_THEME is not None and _THEME.id != "stellar_aiz") else "オレンジ字"
            subtitle_note = f"{color_word}：{continuous_label}（{_num_to_kanji(n_periods)}期）"

    if subtitle_note:
        add_textbox(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=NOTE_HEIGHT,
            text=subtitle_note,
            font_size=font_sizes["note"],
            font_bold=True,
            font_color=COLOR_HIGHLIGHT,
            alignment=PP_ALIGN.RIGHT,
        )
        print(f"  ✓ 注記: {subtitle_note}")

    # レイアウト計算
    note_offset = NOTE_HEIGHT if subtitle_note else Inches(0)
    period_header_top = CONTENT_TOP + note_offset
    table_top = period_header_top + PERIOD_HEADER_HEIGHT

    available_width = CONTENT_WIDTH
    total_gaps = TABLE_GAP * (n_periods - 1)
    table_width = int((available_width - total_gaps) / n_periods)
    table_height = CONTENT_BOTTOM - table_top - BOTTOM_MARGIN

    print(f"  テーブル幅: {table_width/914400:.2f}\" × {n_periods}期")
    print(f"  テーブル高さ: {table_height/914400:.2f}\"")

    for pi, period in enumerate(periods):
        tbl_left = CONTENT_LEFT + (table_width + TABLE_GAP) * pi

        period_label = period.get("label", f"期{pi+1}")
        add_textbox(
            slide,
            left=tbl_left,
            top=period_header_top,
            width=table_width,
            height=PERIOD_HEADER_HEIGHT,
            text=period_label,
            font_size=font_sizes["period"],
            font_bold=True,
            font_color=COLOR_TEXT,
            alignment=PP_ALIGN.CENTER,
            font_underline=True,
        )
        print(f"  ✓ 期ヘッダー: {period_label}")

        build_period_table(
            slide, period, continuous_names, font_sizes,
            tbl_left=tbl_left,
            tbl_top=table_top,
            tbl_width=table_width,
            tbl_height=table_height,
        )
        print(f"  ✓ テーブル[{period_label}]: {len(period.get('customers', []))+1}行 × 4列")

    # Source
    source_text = data.get("source", "")
    if source_text:
        body = source_text if source_text.startswith("出典") else f"出典：{source_text}"
        add_source_label(slide, body)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
