#!/usr/bin/env python3
"""
merge_pptx_v2.py - Robust PPTX merger (v2)

Fixes critical bugs in v1:
  1. Chart _rels files (ppt/charts/_rels/) are properly copied
  2. Excel embedding files (.xlsx) are included in merged output
  3. Content-Type defaults (.png, .xlsx, .bin, etc.) are registered
  4. Chart -> Excel -> Slide reference chain is fully maintained

Usage:
    python merge_pptx_v2.py <output.pptx> <input1.pptx> <input2.pptx> [...]

Architecture:
  - Reads all PPTX files as in-memory ZIP dicts {arcname: bytes}
  - First file becomes the base
  - Subsequent files: all resources are renumbered with offsets to avoid collisions
  - All rels files are updated to reflect renamed targets
  - Content-Types are rebuilt to cover all file extensions
"""

import zipfile
import os
import re
import sys
import shutil
from lxml import etree

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

NSMAP = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
}

SLIDE_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
LAYOUT_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
MASTER_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster"

# Default content types by file extension
EXT_CONTENT_TYPES = {
    "xml": "application/xml",
    "rels": "application/vnd.openxmlformats-package.relationships+xml",
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "emf": "image/x-emf",
    "wmf": "image/x-wmf",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "bin": "application/vnd.openxmlformats-officedocument.oleObject",
    "vml": "application/vnd.openxmlformats-officedocument.vmlDrawing",
    "tiff": "image/tiff",
    "svg": "image/svg+xml",
    "bmp": "image/bmp",
}

# Content types for specific part types (by path pattern)
PART_CONTENT_TYPES = {
    "ppt/slides/slide": "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
    "ppt/slideLayouts/slideLayout": "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml",
    "ppt/slideMasters/slideMaster": "application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml",
    "ppt/theme/theme": "application/vnd.openxmlformats-officedocument.theme+xml",
    "ppt/charts/chart": "application/vnd.openxmlformats-officedocument.drawingml.chart+xml",
    "ppt/tags/tag": "application/vnd.openxmlformats-officedocument.presentationml.tags+xml",
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def ns(prefix, tag):
    return f"{{{NSMAP[prefix]}}}{tag}"


def read_zip(path):
    """Read all entries from a PPTX/ZIP into {arcname: bytes}."""
    contents = {}
    with zipfile.ZipFile(path, 'r') as z:
        for name in z.namelist():
            contents[name] = z.read(name)
    return contents


def write_zip(contents, path):
    """Write {arcname: bytes} to a PPTX/ZIP file."""
    with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as z:
        for name in sorted(contents.keys()):
            z.writestr(name, contents[name])


def max_num(names, pattern):
    """Find the maximum number matching a regex pattern across file names."""
    mx = 0
    for n in names:
        m = re.search(pattern, n)
        if m:
            mx = max(mx, int(m.group(1)))
    return mx


def max_rid(rels_data):
    """Get the highest rId number from a rels XML bytes."""
    root = etree.fromstring(rels_data)
    mx = 0
    for r in root.findall(f"{{{NSMAP['rel']}}}Relationship"):
        m = re.match(r'rId(\d+)', r.get('Id', ''))
        if m:
            mx = max(mx, int(m.group(1)))
    return mx


# ---------------------------------------------------------------------------
# Rename map builder
# ---------------------------------------------------------------------------

def build_rename_map(src, base_names):
    """
    Build a mapping of {src_path -> new_path} for all resources in src
    that need to be copied into base with renumbered filenames.
    
    Also maps corresponding _rels files.
    """
    rmap = {}

    # Resource patterns: (regex matching src paths, format template for new path)
    # Order matters: themes before masters before layouts (dependency order)
    numbered_patterns = [
        (r'ppt/theme/theme(\d+)\.xml$',             'ppt/theme/theme{}.xml'),
        (r'ppt/slideMasters/slideMaster(\d+)\.xml$', 'ppt/slideMasters/slideMaster{}.xml'),
        (r'ppt/slideLayouts/slideLayout(\d+)\.xml$',  'ppt/slideLayouts/slideLayout{}.xml'),
        (r'ppt/slides/slide(\d+)\.xml$',              'ppt/slides/slide{}.xml'),
        (r'ppt/charts/chart(\d+)\.xml$',              'ppt/charts/chart{}.xml'),
        (r'ppt/tags/tag(\d+)\.xml$',                  'ppt/tags/tag{}.xml'),
        (r'ppt/embeddings/oleObject(\d+)\.bin$',      'ppt/embeddings/oleObject{}.bin'),
    ]

    for pattern, template in numbered_patterns:
        # Find matching files in src
        matches = []
        for sp in src:
            m = re.match(pattern, sp)
            if m:
                matches.append((sp, int(m.group(1))))
        matches.sort(key=lambda x: x[1])

        # Find max number in base for this pattern
        base_max = max_num(base_names, pattern)
        counter = base_max

        for src_path, _ in matches:
            counter += 1
            new_path = template.format(counter)
            rmap[src_path] = new_path

            # Also map the corresponding rels file
            src_dir = os.path.dirname(src_path)
            src_file = os.path.basename(src_path)
            old_rels = f"{src_dir}/_rels/{src_file}.rels"
            new_file = os.path.basename(new_path)
            new_rels = f"{os.path.dirname(new_path)}/_rels/{new_file}.rels"
            if old_rels in src:
                rmap[old_rels] = new_rels

    # Handle media files (extension varies: .png, .emf, .jpg, etc.)
    media_max = max_num(base_names, r'ppt/media/image(\d+)')
    counter = media_max
    for sp in sorted(src):
        m = re.match(r'ppt/media/image(\d+)\.(\w+)$', sp)
        if m:
            counter += 1
            new_path = f'ppt/media/image{counter}.{m.group(2)}'
            rmap[sp] = new_path

    # Handle Excel embeddings (various naming patterns)
    excel_max = max_num(base_names, r'Microsoft_Excel_Sheet(\d+)')
    if excel_max == 0:
        excel_max = max_num(base_names, r'Microsoft_Excel_Sheet(\d+)')
    counter = excel_max
    for sp in sorted(src):
        if sp.startswith('ppt/embeddings/') and sp.endswith('.xlsx') and sp not in rmap:
            counter += 1
            rmap[sp] = f'ppt/embeddings/Microsoft_Excel_Sheet{counter}.xlsx'

    return rmap


# ---------------------------------------------------------------------------
# Rels updater
# ---------------------------------------------------------------------------

def apply_rename_to_rels(rels_data, rmap):
    """
    Update all Target attributes in a rels XML according to the rename map.
    Matches by basename to handle relative paths.
    """
    root = etree.fromstring(rels_data)
    for rel in root.findall(f"{{{NSMAP['rel']}}}Relationship"):
        target = rel.get('Target', '')
        basename = os.path.basename(target)
        for old_path, new_path in rmap.items():
            old_base = os.path.basename(old_path)
            if old_base == basename and not old_path.endswith('.rels'):
                new_base = os.path.basename(new_path)
                if new_base != basename:
                    rel.set('Target', target.replace(basename, new_base))
                break
    return etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)


# ---------------------------------------------------------------------------
# Content-Types fixer
# ---------------------------------------------------------------------------

def fix_content_types(base):
    """Ensure all file extensions and parts have proper content type entries."""
    ct_root = etree.fromstring(base['[Content_Types].xml'])

    # Collect existing defaults
    existing_ext = set()
    for el in ct_root.findall(ns('ct', 'Default')):
        existing_ext.add(el.get('Extension'))

    # Collect existing overrides
    existing_parts = set()
    for el in ct_root.findall(ns('ct', 'Override')):
        existing_parts.add(el.get('PartName'))

    # Add missing extension defaults
    for filepath in base:
        if filepath == '[Content_Types].xml':
            continue
        ext = os.path.splitext(filepath)[1].lstrip('.').lower()
        if ext and ext not in existing_ext and ext in EXT_CONTENT_TYPES:
            d = etree.SubElement(ct_root, ns('ct', 'Default'))
            d.set('Extension', ext)
            d.set('ContentType', EXT_CONTENT_TYPES[ext])
            existing_ext.add(ext)

    # Add missing part overrides
    for filepath in base:
        partname = '/' + filepath
        if partname in existing_parts:
            continue
        for prefix, ctype in PART_CONTENT_TYPES.items():
            if filepath.startswith(prefix) and filepath.endswith('.xml'):
                o = etree.SubElement(ct_root, ns('ct', 'Override'))
                o.set('PartName', partname)
                o.set('ContentType', ctype)
                existing_parts.add(partname)
                break

    base['[Content_Types].xml'] = etree.tostring(
        ct_root, xml_declaration=True, encoding='UTF-8', standalone=True
    )


# ---------------------------------------------------------------------------
# Main merge logic
# ---------------------------------------------------------------------------

def _finalize_pptx(path):
    """LibreOffice headless roundtrip to normalize OOXML and clear PowerPoint
    "修復が必要" ダイアログ triggers. No-op if soffice is unavailable or times out;
    the original file is kept on any failure."""
    import glob
    import subprocess
    import tempfile
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=180, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


def merge_presentations(input_files, output_path, *, roundtrip=True):
    """Merge multiple PPTX files into one output file."""

    if not input_files:
        print("Error: No input files provided.")
        sys.exit(1)

    if len(input_files) == 1:
        shutil.copy2(input_files[0], output_path)
        if roundtrip:
            _finalize_pptx(output_path)
        print(f"Only one file provided. Copied to {output_path}")
        return

    # Read the first file as base
    base = read_zip(input_files[0])

    for src_idx, src_file in enumerate(input_files[1:], 2):
        print(f"Merging: {os.path.basename(src_file)} ({src_idx}/{len(input_files)})")
        src = read_zip(src_file)

        # Build rename map for all resources
        rmap = build_rename_map(src, list(base.keys()))

        # Copy all renamed resources into base
        for old_path, new_path in rmap.items():
            if old_path in src:
                data = src[old_path]
                # If it's a rels file, update targets according to rename map
                if old_path.endswith('.rels'):
                    try:
                        data = apply_rename_to_rels(data, rmap)
                    except Exception as e:
                        print(f"  Warning: Could not update rels {old_path}: {e}")
                base[new_path] = data

        # --- Update presentation.xml: add masters and slides ---
        base_pres = etree.fromstring(base['ppt/presentation.xml'])
        base_prels = etree.fromstring(base['ppt/_rels/presentation.xml.rels'])
        rid_counter = max_rid(base['ppt/_rels/presentation.xml.rels'])

        # Add slide masters
        master_id_lst = base_pres.find(f".//{ns('p', 'sldMasterIdLst')}")
        if master_id_lst is None:
            master_id_lst = etree.SubElement(base_pres, ns('p', 'sldMasterIdLst'))

        max_mid = 2147483647
        for el in master_id_lst.findall(ns('p', 'sldMasterId')):
            max_mid = max(max_mid, int(el.get('id', '0')))
        max_mid = max(max_mid, 2147483647)

        for old_p, new_p in sorted(rmap.items()):
            if re.match(r'ppt/slideMasters/slideMaster\d+\.xml$', new_p):
                rid_counter += 1
                max_mid += 1
                rid = f'rId{rid_counter}'

                r = etree.SubElement(base_prels, f"{{{NSMAP['rel']}}}Relationship")
                r.set('Id', rid)
                r.set('Type', MASTER_RT)
                r.set('Target', new_p.replace('ppt/', ''))

                m = etree.SubElement(master_id_lst, ns('p', 'sldMasterId'))
                m.set('id', str(max_mid))
                m.set(ns('r', 'id'), rid)

        # Get slide order from source presentation.xml
        src_pres = etree.fromstring(src['ppt/presentation.xml'])
        src_prels = etree.fromstring(src['ppt/_rels/presentation.xml.rels'])

        src_rid_slide = {}
        for r in src_prels.findall(f"{{{NSMAP['rel']}}}Relationship"):
            if r.get('Type') == SLIDE_RT:
                src_rid_slide[r.get('Id')] = os.path.basename(r.get('Target'))

        ordered = []
        src_sld_lst = src_pres.find(f".//{ns('p', 'sldIdLst')}")
        if src_sld_lst is not None:
            for sid in src_sld_lst.findall(ns('p', 'sldId')):
                rid = sid.get(ns('r', 'id'))
                if rid in src_rid_slide:
                    ordered.append(src_rid_slide[rid])

        # Add slides to base
        sld_lst = base_pres.find(f".//{ns('p', 'sldIdLst')}")
        if sld_lst is None:
            sld_lst = etree.SubElement(base_pres, ns('p', 'sldIdLst'))

        max_sid = 255
        for s in sld_lst.findall(ns('p', 'sldId')):
            max_sid = max(max_sid, int(s.get('id', '0')))

        for slide_file in ordered:
            old_slide = f'ppt/slides/{slide_file}'
            if old_slide not in rmap:
                continue
            new_slide = rmap[old_slide]

            rid_counter += 1
            max_sid += 1
            rid = f'rId{rid_counter}'

            r = etree.SubElement(base_prels, f"{{{NSMAP['rel']}}}Relationship")
            r.set('Id', rid)
            r.set('Type', SLIDE_RT)
            r.set('Target', new_slide.replace('ppt/', ''))

            s = etree.SubElement(sld_lst, ns('p', 'sldId'))
            s.set('id', str(max_sid))
            s.set(ns('r', 'id'), rid)

        # Save updated presentation XML
        base['ppt/presentation.xml'] = etree.tostring(
            base_pres, xml_declaration=True, encoding='UTF-8', standalone=True
        )
        base['ppt/_rels/presentation.xml.rels'] = etree.tostring(
            base_prels, xml_declaration=True, encoding='UTF-8', standalone=True
        )

    # Fix content types for all extensions and parts
    fix_content_types(base)

    # Write output
    write_zip(base, output_path)

    # --- Verification ---
    print(f"\n{'='*60}")
    print(f"Verification")
    print(f"{'='*60}")

    try:
        from pptx import Presentation
        prs = Presentation(output_path)
        print(f"Slides: {len(prs.slides)}")
        for i, sl in enumerate(prs.slides, 1):
            txt = ''
            for sh in sl.shapes:
                if hasattr(sh, 'text') and sh.text.strip() and len(sh.text.strip()) > 10:
                    txt = sh.text.strip()[:60]
                    break
            print(f"  Slide {i}: {len(sl.shapes)} shapes | {txt}")
    except ImportError:
        print("  (python-pptx not available for verification)")
    except Exception as e:
        print(f"  Verification warning: {e}")

    # Check chart -> Excel chain
    with zipfile.ZipFile(output_path) as z:
        names = z.namelist()
        charts = [f for f in names if re.match(r'ppt/charts/chart\d+\.xml$', f)]
        for chart in charts:
            rels_path = chart.replace('charts/', 'charts/_rels/') + '.rels'
            if rels_path in names:
                root = etree.fromstring(z.read(rels_path))
                for rel in root.findall(f"{{{NSMAP['rel']}}}Relationship"):
                    tgt = rel.get('Target', '')
                    resolved = os.path.normpath(f"ppt/charts/{tgt}").replace('\\', '/')
                    if '../' in tgt:
                        resolved = resolved.replace('ppt/charts/../', 'ppt/')
                    exists = resolved in names
                    status = "OK" if exists else "MISSING"
                    if not exists:
                        print(f"  WARNING: {chart} -> {tgt} [{status}]")
            else:
                print(f"  WARNING: No rels for {chart}")

    # OOXML normalization via LibreOffice (clears "修復が必要" on most outputs).
    if roundtrip:
        _finalize_pptx(output_path)

    print(f"\nMerged {len(input_files)} files -> {output_path}")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    args = sys.argv[1:]
    roundtrip = True
    if "--no-roundtrip" in args:
        roundtrip = False
        args = [a for a in args if a != "--no-roundtrip"]

    if len(args) < 2:
        print("Usage: python merge_pptx_v2.py [--no-roundtrip] <output.pptx> <input1.pptx> <input2.pptx> [...]")
        sys.exit(1)

    output = args[0]
    inputs = args[1:]

    for f in inputs:
        if not os.path.exists(f):
            print(f"Error: File not found: {f}")
            sys.exit(1)

    merge_presentations(inputs, output, roundtrip=roundtrip)
