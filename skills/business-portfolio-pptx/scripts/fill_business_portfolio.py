"""
fill_business_portfolio.py — 事業ポートフォリオスライドをPPTXネイティブオブジェクトで生成

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 左側: セグメント別売上高の積み上げ棒グラフ (絶対値、複数年)
  - 右側: セグメント別サマリーテーブル (最新期売上、構成比、CAGR、営業利益率)
  - 下部: 出典 (stella=動的 textbox / roleup=Source 3 placeholder)

Roleup 固有:
  - 茶色ガイド `正方形/長方形 1` `正方形/長方形 8` を silent_remove
  - chart.has_legend = False、custom legend を chart 下に配置 (C12 対策)
  - data label / cat axis = 10pt、source = 6pt
  - segment 色未指定時は theme.chart_palette を使用

Usage:
  python fill_business_portfolio.py --brand stellar_aiz \\
    --data {{WORK_DIR}}/business_portfolio_data.json \\
    --output {{OUTPUT_DIR}}/BusinessPortfolio_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "business-portfolio-pptx"

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

# Defaults (stella). Reassigned in _apply_theme(theme) for roleup.
SHAPE_SOURCE = "Source"

PANEL_Y = Inches(1.50)
LEFT_X = Inches(0.41)
LEFT_W = Inches(6.30)
LEFT_H = Inches(5.20)
RIGHT_X = Inches(7.00)
RIGHT_W = Inches(5.90)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)
SOURCE_H = Inches(0.30)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_HEADER_BG = RGBColor(0x2E, 0x4A, 0x6B)
COLOR_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_ALT = RGBColor(0xF2, 0xF2, 0xF2)

DEFAULT_COLORS = [
    "#4E79A7", "#F28E2B", "#59A14F", "#E15759", "#76B7B2",
    "#EDC948", "#B07AA1", "#FF9DA7", "#9C755F", "#BAB0AC",
]

FONT_NAME_JP = "Meiryo UI"
TEXT_HEX = "333333"
HEADER_BG_HEX = "2E4A6B"
HEADER_TEXT_HEX = "FFFFFF"
ROW_ALT_HEX = "F2F2F2"

FONT_SIZE_SECTION_PT = 14
FONT_SIZE_SECTION_BOLD = True
SECTION_TITLE_HEX = "333333"
SECTION_TITLE_ALIGN = PP_ALIGN.CENTER

FONT_SIZE_TABLE_PT = 11
FONT_SIZE_TABLE_HEADER_PT = 11
FONT_SIZE_DATA_LABEL_PT = 9
FONT_SIZE_CAT_AXIS_PT = 11
FONT_SIZE_LEGEND_PT = 10
FONT_SIZE_SOURCE_PT = 10
FONT_SIZE_UNIT_PT = 11

USE_BUILTIN_LEGEND = True

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global SHAPE_SOURCE
    global PANEL_Y, LEFT_X, LEFT_W, LEFT_H, RIGHT_X, RIGHT_W
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE, COLOR_HEADER_BG, COLOR_HEADER_TEXT, COLOR_ROW_ALT
    global DEFAULT_COLORS, FONT_NAME_JP, TEXT_HEX, HEADER_BG_HEX, HEADER_TEXT_HEX, ROW_ALT_HEX
    global FONT_SIZE_SECTION_PT, FONT_SIZE_SECTION_BOLD, SECTION_TITLE_HEX, SECTION_TITLE_ALIGN
    global FONT_SIZE_TABLE_PT, FONT_SIZE_TABLE_HEADER_PT
    global FONT_SIZE_DATA_LABEL_PT, FONT_SIZE_CAT_AXIS_PT, FONT_SIZE_LEGEND_PT
    global FONT_SIZE_SOURCE_PT, FONT_SIZE_UNIT_PT, USE_BUILTIN_LEGEND

    _THEME = theme
    FONT_NAME_JP = theme.font_ea
    TEXT_HEX = theme.hex_no_hash("text")

    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    PANEL_Y = theme.layout("panel_y_in")
    LEFT_X = theme.layout("left_x_in")
    LEFT_W = theme.layout("left_w_in")
    LEFT_H = theme.layout("left_h_in")
    RIGHT_X = theme.layout("right_x_in")
    RIGHT_W = theme.layout("right_w_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    DEFAULT_COLORS = list(theme.chart_palette)

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        COLOR_HEADER_BG = RGBColor(0x2E, 0x4A, 0x6B)
        COLOR_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
        COLOR_ROW_ALT = RGBColor(0xF2, 0xF2, 0xF2)
        HEADER_BG_HEX = "2E4A6B"
        HEADER_TEXT_HEX = "FFFFFF"
        ROW_ALT_HEX = "F2F2F2"
        FONT_SIZE_SECTION_PT = 14
        FONT_SIZE_SECTION_BOLD = True
        SECTION_TITLE_HEX = TEXT_HEX
        SECTION_TITLE_ALIGN = PP_ALIGN.CENTER
        FONT_SIZE_TABLE_PT = 11
        FONT_SIZE_TABLE_HEADER_PT = 11
        FONT_SIZE_DATA_LABEL_PT = 9
        FONT_SIZE_CAT_AXIS_PT = 11
        FONT_SIZE_LEGEND_PT = 10
        FONT_SIZE_SOURCE_PT = 10
        FONT_SIZE_UNIT_PT = 11
        USE_BUILTIN_LEGEND = True
    else:
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        SHAPE_SOURCE = "Source 3"
        # Header: light beige (header_bg) + text color (dark) → C4 brand-aligned
        COLOR_HEADER_BG = theme.color("header_bg")
        COLOR_HEADER_TEXT = theme.color("text")
        COLOR_ROW_ALT = theme.color("label_bg")
        HEADER_BG_HEX = theme.hex_no_hash("header_bg")
        HEADER_TEXT_HEX = theme.hex_no_hash("text")
        ROW_ALT_HEX = theme.hex_no_hash("label_bg")
        FONT_SIZE_SECTION_PT = theme.pt_value("font_size_subtitle_pt")  # 12
        FONT_SIZE_SECTION_BOLD = False
        SECTION_TITLE_HEX = theme.hex_no_hash("subtitle")  # #897141
        SECTION_TITLE_ALIGN = PP_ALIGN.LEFT
        FONT_SIZE_TABLE_PT = theme.pt_value("font_size_body_pt")  # 10
        FONT_SIZE_TABLE_HEADER_PT = theme.pt_value("font_size_body_pt")  # 10
        FONT_SIZE_DATA_LABEL_PT = theme.pt_value("font_size_body_pt")  # 10 (was 9)
        FONT_SIZE_CAT_AXIS_PT = theme.pt_value("font_size_body_pt")  # 10
        FONT_SIZE_LEGEND_PT = theme.pt_value("font_size_body_pt")  # 10 (custom legend)
        FONT_SIZE_SOURCE_PT = theme.pt_value("font_size_source_pt")  # 6
        FONT_SIZE_UNIT_PT = theme.pt_value("font_size_body_pt")  # 10
        USE_BUILTIN_LEGEND = False


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_shape(slide, name, warn=True):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    if warn:
        print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def write_source_placeholder(shape, text, font_size_pt):
    """Source 3 placeholder にテキストを書き込み、フォント size/name を強制。"""
    tf = shape.text_frame
    para = tf.paragraphs[0]
    for r in list(para.runs):
        r.text = ""
    if para.runs:
        run = para.runs[0]
        run.text = text
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text
        run = para.runs[0]
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        run._r.insert(0, rPr)
    rPr.set("sz", str(font_size_pt * 100))
    for tag in [qn("a:latin"), qn("a:ea")]:
        old = rPr.find(tag)
        if old is not None:
            rPr.remove(old)
        etree.SubElement(rPr, tag, attrib={"typeface": FONT_NAME_JP})


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_section_title(slide, text, left, top, width):
    """セクションタイトル (Bold + 下線) を追加。"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = SECTION_TITLE_ALIGN
    run = p.add_run()
    run.text = text
    run.font.size = Pt(FONT_SIZE_SECTION_PT)
    run.font.bold = FONT_SIZE_SECTION_BOLD
    run.font.color.rgb = hex_to_rgb(SECTION_TITLE_HEX)
    run.font.name = FONT_NAME_JP

    # 下線 (装飾用シェイプ、テキストなしのため C11 対象外)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(SECTION_TITLE_HEX)
    line.line.fill.background()
    return txBox


def add_text_box(slide, text, left, top, width, height, font_size_pt, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.name = FONT_NAME_JP
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


def add_custom_legend(slide, segments, left, top, width):
    """Custom legend (C12 対策、roleup)。■ + segment name を横並びで配置。"""
    legend_h = Inches(0.25)
    sq_size = Inches(0.14)
    gap = Inches(0.10)
    item_gap = Inches(0.30)

    # 各 item: ■ + " name " 幅見積。日本語名は概算 0.8in/item に固定。
    item_w = (width - item_gap * (len(segments) - 1)) // max(1, len(segments))
    if item_w < Inches(1.0):
        item_w = Inches(1.0)

    cur_x = left
    for seg in segments:
        # ■ marker
        sq_y = top + (legend_h - sq_size) // 2
        sq = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(cur_x), int(sq_y), sq_size, sq_size,
        )
        sq.fill.solid()
        sq.fill.fore_color.rgb = hex_to_rgb(seg["color"])
        sq.line.fill.background()

        text_x = cur_x + sq_size + gap
        text_w = item_w - sq_size - gap
        tb = slide.shapes.add_textbox(int(text_x), top, text_w, legend_h)
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = 0; tf.margin_right = 0
        tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = seg["name"]
        run.font.size = Pt(FONT_SIZE_LEGEND_PT)
        run.font.color.rgb = COLOR_TEXT
        run.font.name = FONT_NAME_JP

        cur_x += item_w + item_gap


def build_stacked_bar_chart(slide, section_title, chart_data, segments, left, top, width, height):
    add_section_title(slide, section_title, left, top, width)

    unit_label = chart_data.get("unit_label", "（単位：億円）")
    add_text_box(
        slide, unit_label,
        left, top + Inches(0.40), Inches(2.0), Inches(0.25),
        FONT_SIZE_UNIT_PT, bold=False, align=PP_ALIGN.LEFT,
    )

    # チャート位置: タイトル + 単位 の下、底部に custom legend スペースを残す
    chart_top = top + Inches(0.70)
    legend_reserve = Inches(0.0) if USE_BUILTIN_LEGEND else Inches(0.40)
    chart_h = height - Inches(0.70) - legend_reserve

    years = chart_data["years"]
    values = chart_data["values"]

    cdata = CategoryChartData()
    cdata.categories = years

    seg_names = [seg["name"] for seg in segments]
    for seg_name in seg_names:
        series_vals = values.get(seg_name, [0] * len(years))
        cdata.add_series(seg_name, series_vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        left, chart_top, width, chart_h,
        cdata,
    )
    chart = chart_shape.chart

    if USE_BUILTIN_LEGEND:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(FONT_SIZE_LEGEND_PT)
        chart.legend.font.name = FONT_NAME_JP
    else:
        chart.has_legend = False

    chart.has_title = False

    for idx, series in enumerate(chart.series):
        seg = segments[idx] if idx < len(segments) else {}
        hex_color = seg.get("color") or DEFAULT_COLORS[idx % len(DEFAULT_COLORS)]
        rgb = hex_to_rgb(hex_color)

        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = rgb

        series.format.line.fill.background()

        series.data_labels.show_value = True
        series.data_labels.position = XL_LABEL_POSITION.CENTER
        series.data_labels.font.size = Pt(FONT_SIZE_DATA_LABEL_PT)
        series.data_labels.font.name = FONT_NAME_JP
        series.data_labels.font.bold = True
        series.data_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        series.data_labels.number_format = '#,##0'
        series.data_labels.number_format_is_linked = False

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(FONT_SIZE_CAT_AXIS_PT)
    cat_axis.tick_labels.font.name = FONT_NAME_JP

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(FONT_SIZE_CAT_AXIS_PT)
    val_axis.tick_labels.font.name = FONT_NAME_JP
    val_axis.visible = False
    val_axis.major_tick_mark = 2

    year_totals = []
    for y_idx in range(len(years)):
        total = sum(values.get(seg_name, [0] * len(years))[y_idx] for seg_name in seg_names)
        year_totals.append(total)
    max_total = max(year_totals) if year_totals else 0
    if max_total > 0:
        val_axis.maximum_scale = max_total * 1.20
        val_axis.minimum_scale = 0

    # Custom legend (roleup) — chart bottom の外側に配置
    if not USE_BUILTIN_LEGEND:
        legend_y = chart_top + chart_h + Inches(0.05)
        add_custom_legend(slide, segments, left, legend_y, width)

    print(f"  ✓ 積み上げ棒チャート: {len(seg_names)}セグメント x {len(years)}年")
    return chart_shape


def _header_to_key(header, col_idx):
    mapping = {
        "セグメント": "name",
        "事業セグメント": "name",
        "セグメント名": "name",
        "事業": "name",
        "売上": "revenue",
        "売上高": "revenue",
        "最新期売上": "revenue",
        "構成比": "share",
        "シェア": "share",
        "CAGR": "cagr",
        "成長率": "cagr",
        "利益率": "op_margin",
        "営業利益率": "op_margin",
        "営業利益": "op_profit",
    }
    return mapping.get(header, f"col{col_idx}")


def _style_cell(cell, text, is_header=False, is_alt=False, font_size_pt=11,
                bold=False, color_marker_hex=None):
    tc = cell._tc

    if is_header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
    elif is_alt:
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ROW_ALT
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = cell.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    p_elem = etree.SubElement(tf._txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "ctr" if is_header else ("l" if bold else "r"))
    if not bold and not is_header:
        pPr.set("algn", "r")

    text_color_hex = HEADER_TEXT_HEX if is_header else TEXT_HEX

    if color_marker_hex:
        # 2 runs: ■ (segment color) + name (text color)
        marker_rgb = color_marker_hex.lstrip("#")
        r1 = etree.SubElement(p_elem, qn("a:r"))
        rPr1 = etree.SubElement(r1, qn("a:rPr"), attrib={
            "lang": "ja-JP", "sz": str(font_size_pt * 100)
        })
        rPr1.set("b", "1")
        etree.SubElement(rPr1, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr1, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf1 = etree.SubElement(rPr1, qn("a:solidFill"))
        s1 = etree.SubElement(sf1, qn("a:srgbClr"))
        s1.set("val", marker_rgb)
        t1 = etree.SubElement(r1, qn("a:t"))
        t1.text = "■ "

        r2 = etree.SubElement(p_elem, qn("a:r"))
        rPr2 = etree.SubElement(r2, qn("a:rPr"), attrib={
            "lang": "ja-JP", "sz": str(font_size_pt * 100)
        })
        rPr2.set("b", "1")
        etree.SubElement(rPr2, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr2, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf2 = etree.SubElement(rPr2, qn("a:solidFill"))
        s2 = etree.SubElement(sf2, qn("a:srgbClr"))
        s2.set("val", text_color_hex)
        t2 = etree.SubElement(r2, qn("a:t"))
        t2.text = text
    else:
        r_elem = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
            "lang": "ja-JP", "sz": str(font_size_pt * 100)
        })
        rPr.set("b", "1" if (is_header or bold) else "0")
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        sr = etree.SubElement(sf, qn("a:srgbClr"))
        sr.set("val", text_color_hex)
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def build_summary_table(slide, table_data, segments, left, top, width):
    section_title = table_data.get("section_title", "セグメント別サマリー")
    add_section_title(slide, section_title, left, top, width)

    headers = table_data["headers"]
    rows = table_data["rows"]

    n_cols = len(headers)
    n_rows = len(rows) + 1

    tbl_top = top + Inches(0.55)
    row_h = Inches(0.38)
    tbl_h = row_h * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, tbl_top, width, tbl_h)
    table = shape.table

    first_col_w = int(width * 0.32)
    other_col_w = int((width - first_col_w) / (n_cols - 1))
    table.columns[0].width = Emu(first_col_w)
    for i in range(1, n_cols):
        table.columns[i].width = Emu(other_col_w)

    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    for tr in tbl_elem.findall(qn('a:tr')):
        tr.set('h', str(row_h))

    for c_idx, h in enumerate(headers):
        cell = table.cell(0, c_idx)
        _style_cell(cell, h, is_header=True, is_alt=False,
                    font_size_pt=FONT_SIZE_TABLE_HEADER_PT)

    seg_color_map = {seg["name"]: seg.get("color") for seg in segments}
    for r_idx, row in enumerate(rows):
        is_alt = (r_idx % 2 == 1)
        for c_idx, h in enumerate(headers):
            key = _header_to_key(h, c_idx)
            val = row.get(key, "")
            if c_idx == 0:
                cell = table.cell(r_idx + 1, c_idx)
                _style_cell(cell, val, is_header=False, is_alt=is_alt,
                            font_size_pt=FONT_SIZE_TABLE_PT, bold=True,
                            color_marker_hex=seg_color_map.get(val))
            else:
                cell = table.cell(r_idx + 1, c_idx)
                _style_cell(cell, val, is_header=False, is_alt=is_alt,
                            font_size_pt=FONT_SIZE_TABLE_PT)

    print(f"  ✓ サマリーテーブル: {n_rows}行 x {n_cols}列")
    return shape


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True, help="Path to JSON data file")
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True, help="Path to output PPTX")
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "business-portfolio")
    print(f"=== 事業ポートフォリオスライド生成 (brand={theme.id}) ===")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    require_source(data, theme, skill_id=SKILL_ID)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Roleup: silently remove brown guide rectangles
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # Top placeholder (stella: main_message / roleup: chart_title)
    top_text = resolve_top_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:50]}")

    # Subtitle placeholder
    sub_text = resolve_subtitle_text(data, theme) or "事業ポートフォリオ"
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:50]}")

    segments = data.get("segments", [])
    if not segments:
        print("  ✗ ERROR: 'segments' is required", file=sys.stderr)
        sys.exit(1)

    for i, seg in enumerate(segments):
        if not seg.get("color"):
            seg["color"] = DEFAULT_COLORS[i % len(DEFAULT_COLORS)]

    chart_data = data.get("chart", {})
    left_section_title = chart_data.get("section_title", "セグメント別売上高推移")
    build_stacked_bar_chart(
        slide, left_section_title, chart_data, segments,
        LEFT_X, PANEL_Y, LEFT_W, LEFT_H,
    )

    table_data = data.get("table", {})
    if table_data:
        build_summary_table(slide, table_data, segments, RIGHT_X, PANEL_Y, RIGHT_W)

    source = data.get("source", "")
    if source:
        body = source if source.startswith("出典") else f"出典：{source}"
        source_shape = find_shape(slide, SHAPE_SOURCE, warn=False)
        if source_shape is not None:
            write_source_placeholder(source_shape, body, FONT_SIZE_SOURCE_PT)
            print(f"  ✓ 出典 ({SHAPE_SOURCE} placeholder): {body[:50]}")
        else:
            add_text_box(
                slide, body,
                SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                FONT_SIZE_SOURCE_PT, bold=False, color=COLOR_SOURCE,
                align=PP_ALIGN.LEFT,
            )
            print(f"  ✓ 出典 (dynamic textbox): {body[:50]}")

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
